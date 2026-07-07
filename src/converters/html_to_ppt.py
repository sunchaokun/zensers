# -*- coding: utf-8 -*-
"""
HTML to PPT Converter
=====================

Convert HTML intermediate format to PowerPoint document (.pptx).

Core Features:
1. HTML parsing and slide conversion
2. Slide type handling (cover, toc, content, data, end)
3. Title and content layout
4. List and table processing
5. Style application

Technology: python-pptx

Usage Example:
    converter = HTMLToPPTConverter()
    result = converter.convert(
        html="<section class='slide' data-type='cover'><h1>Title</h1></section>",
        output_path="output.pptx"
    )
"""

import logging
import os
import re
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

from .base_parser import SlideElementParser
from .css_extractor import CSSStyleExtractor, ExtractedStyles

logger = logging.getLogger(__name__)

# Constants
MAX_HTML_SIZE = 50 * 1024 * 1024  # 50MB


@dataclass
class ConversionResult:
    """Conversion Result"""
    success: bool
    output_path: Optional[str] = None
    file_size: Optional[int] = None
    slides_count: Optional[int] = None
    error: Optional[str] = None
    error_code: Optional[str] = None
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary"""
        return {
            "success": self.success,
            "output_path": self.output_path,
            "file_size": self.file_size,
            "slides_count": self.slides_count,
            "error": self.error,
            "error_code": self.error_code
        }


class HTMLToPPTConverter:
    """
    HTML to PPT Converter
    
    Convert HTML intermediate format to PowerPoint document.
    
    Usage Example:
        converter = HTMLToPPTConverter()
        result = converter.convert(
            html="<section class='slide' data-type='cover'><h1>Title</h1></section>",
            output_path="output.pptx"
        )
        
        if result.success:
            print(f"File saved: {result.output_path}")
    """
    
    # Default styles
    DEFAULT_STYLES = {
        "title_font": "Microsoft YaHei",
        "body_font": "Microsoft YaHei",
        "title_size": 44,
        "subtitle_size": 28,
        "body_size": 18,
        "slide_width": 13.333,  # inches (16:9)
        "slide_height": 7.5  # inches (16:9)
    }
    
    def __init__(self, styles: Optional[Dict[str, Any]] = None):
        """
        Initialize converter
        
        Args:
            styles: Custom style configuration
        """
        self.styles = {**self.DEFAULT_STYLES, **(styles or {})}
        self._pptx_available = self._check_pptx_available()
        
        if not self._pptx_available:
            logger.warning("python-pptx not available, converter will have limited functionality")
    
    def _check_pptx_available(self) -> bool:
        """Check if python-pptx is available"""
        try:
            from pptx import Presentation  # noqa: F401
            return True
        except ImportError:
            return False
    
    def get_default_styles(self) -> Dict[str, Any]:
        """Get default styles"""
        return self.DEFAULT_STYLES.copy()
    
    def _merge_styles(
        self,
        extracted_styles: Optional[ExtractedStyles],
        custom_styles: Optional[Dict[str, Any]]
    ) -> Dict[str, Any]:
        """
        Merge styles
        
        Merge order: template styles → default styles → custom styles
        
        Args:
            extracted_styles: Styles extracted from template
            custom_styles: User custom styles
            
        Returns:
            Merged style dictionary
        """
        # Default styles as base
        final_styles = self.DEFAULT_STYLES.copy()
        
        # Apply template-extracted styles (higher priority than defaults)
        if extracted_styles:
            template_styles = extracted_styles.to_ppt_styles()
            final_styles.update(template_styles)
            logger.debug(f"Applied template styles: {template_styles}")
        
        # Apply custom styles (highest priority)
        if custom_styles:
            final_styles.update(custom_styles)
            logger.debug(f"Applied custom styles: {custom_styles}")
        
        return final_styles
    
    def convert(
        self,
        html: str,
        output_path: str,
        styles: Optional[Dict[str, Any]] = None,
        template_html: Optional[str] = None
    ) -> ConversionResult:
        """
        Convert HTML to PPT document
        
        Args:
            html: HTML content
            output_path: Output file path
            styles: Custom styles (optional)
            template_html: HTML template content (for extracting CSS styles)
            
        Returns:
            Conversion Result
        """
        # Input validation
        if not isinstance(html, str):
            logger.warning("html is not a string, converting to empty string")
            html = ""
        
        if not isinstance(output_path, str):
            return ConversionResult(
                success=False,
                error="output_path must be a string",
                error_code="INVALID_PATH_TYPE"
            )
        
        # File size limit check
        if len(html) > MAX_HTML_SIZE:
            return ConversionResult(
                success=False,
                error=f"HTML content too large ({len(html)/1024/1024:.1f}MB > {MAX_HTML_SIZE/1024/1024}MB)",
                error_code="CONTENT_TOO_LARGE"
            )
        
        # Path security check
        if not self._is_safe_path(output_path):
            return ConversionResult(
                success=False,
                error="Invalid or unsafe output path",
                error_code="UNSAFE_PATH"
            )
        
        # Empty HTML handling
        if not html.strip():
            logger.info("Empty HTML, creating minimal presentation")
            html = "<section class='slide' data-type='cover'><h1></h1></section>"
        
        # Preprocess HTML (cleanup + Markdown conversion)
        html = self._sanitize_html(html)
        
        # Extract CSS styles from template HTML (if any)
        extracted_styles = None
        if template_html:
            extractor = CSSStyleExtractor()
            extracted_styles = extractor.extract_from_html(template_html)
            logger.info("Extracted styles from template HTML")
        
        # Merge styles: template → default → custom
        final_styles = self._merge_styles(extracted_styles, styles)
        
        try:
            # Ensure output directory exists
            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir, exist_ok=True)
            
            # Use slide parser
            parser = SlideElementParser()
            parser.feed(html)
            slides = parser.get_slides()
            
            # Create PPT document
            if self._pptx_available:
                return self._create_pptx_document(slides, output_path, final_styles)
            else:
                return self._create_fallback_document(slides, output_path, final_styles)
                
        except (OSError, IOError) as e:
            logger.error(f"File operation failed: {e}")
            return ConversionResult(
                success=False,
                error=f"File operation failed: {e}",
                error_code="FILE_ERROR"
            )
        except ValueError as e:
            logger.error(f"Invalid data: {e}")
            return ConversionResult(
                success=False,
                error=f"Invalid data: {e}",
                error_code="DATA_ERROR"
            )
        except RuntimeError as e:
            logger.error(f"Processing error: {e}")
            return ConversionResult(
                success=False,
                error=f"Processing error: {e}",
                error_code="PROCESSING_ERROR"
            )
    
    def _is_safe_path(self, path: str, base_dir: Optional[str] = None) -> bool:
        """
        Check if path is safe
        
        Args:
            path: File path
            base_dir: Allowed base directory (optional, defaults to current working directory)
            
        Returns:
            Whether safe
        """
        try:
            resolved_path = Path(path).resolve()
            
            # Check path traversal
            if '..' in path:
                return False
            
            # Determine base directory (default to current working directory if not specified)
            if base_dir is None:
                base_dir = str(Path.cwd())
            
            base_path = Path(base_dir).resolve()
            
            # Check if absolute path is within allowed directory
            try:
                if not resolved_path.is_relative_to(base_path):
                    return False
            except AttributeError:
                # Fallback for Python < 3.9
                try:
                    resolved_path.relative_to(base_path)
                except ValueError:
                    return False
            
            # Block access to sensitive system directories (exact match)
            dangerous_paths = [
                Path('/etc'), Path('/sys'), Path('/proc'), Path('/root'),
                Path('C:/Windows'), Path('C:/Windows/System32'),
                Path('D:/Windows'),
            ]
            for dangerous in dangerous_paths:
                try:
                    if resolved_path.is_relative_to(dangerous):
                        return False
                except (AttributeError, ValueError):
                    if str(resolved_path).startswith(str(dangerous)):
                        return False
            
            return True
            
        except (OSError, ValueError):
            return False
    
    def _sanitize_html(self, html: str) -> str:
        """
        Clean HTML content, remove content that should not appear in PPT
        
        1. Remove <style> tags and content
        2. Remove residual template tags {% %} and {{ }}
        3. Remove other non-content tags
        4. Convert Markdown syntax to HTML tags
        
        Args:
            html: Raw HTML content
            
        Returns:
            Cleaned HTML
        """
        import re
        
        # Remove <style> tags and content
        html = re.sub(r'<style[^>]*>.*?</style>', '', html, flags=re.DOTALL | re.IGNORECASE)
        
        # Remove <script> tags and content
        html = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
        
        # Remove <head> tags and content
        html = re.sub(r'<head[^>]*>.*?</head>', '', html, flags=re.DOTALL | re.IGNORECASE)
        
        # Remove residual template tags
        html = re.sub(r'\{%.*?%\}', '', html, flags=re.DOTALL)
        html = re.sub(r'\{\{.*?\}\}', '', html, flags=re.DOTALL)
        
        # Convert Markdown syntax to HTML tags
        html = self._convert_markdown_to_html(html)
        
        # Clean up excess whitespace
        html = re.sub(r'\n\s*\n\s*\n', '\n\n', html)
        
        return html.strip()
    
    def _convert_markdown_to_html(self, text: str) -> str:
        """
        Convert Markdown syntax to HTML tags
        
        Processed syntax:
        - **bold** → <strong>bold</strong>
        - *italic* → <em>italic</em>
        - `code` → <code>code</code>
        - [link](url) → <a href="url">link</a>
        - ~~strikethrough~~ → <del>strikethrough</del>
        
        Args:
            text: Text that may contain Markdown syntax
            
        Returns:
            Converted HTML
        """
        import re
        
        # Split text into HTML tags and non-tag parts
        parts = re.split(r'(<[^>]+>)', text)
        
        result_parts = []
        for part in parts:
            # If it's an HTML tag, keep unchanged
            if part.startswith('<') and part.endswith('>'):
                result_parts.append(part)
            else:
                # Non-tag part, convert Markdown syntax
                converted = self._convert_markdown_inline(part)
                result_parts.append(converted)
        
        return ''.join(result_parts)
    
    def _convert_markdown_inline(self, text: str) -> str:
        """
        Convert inline Markdown syntax to HTML
        
        Args:
            text: Plain text content
            
        Returns:
            Converted HTML
        """
        import re
        
        # 1. Strikethrough ~~text~~
        text = re.sub(r'~~([^~]+)~~', r'<del>\1</del>', text)
        
        # 2. Bold **text**
        text = re.sub(r'\*\*([^*]+)\*\*', r'<strong>\1</strong>', text)
        
        # 3. Italic *text* or _text_
        text = re.sub(r'\*([^*]+)\*', r'<em>\1</em>', text)
        text = re.sub(r'_([^_]+)_', r'<em>\1</em>', text)
        
        # 4. Inline code `code`
        text = re.sub(r'`([^`]+)`', r'<code>\1</code>', text)
        
        # 5. Link [text](url)
        text = re.sub(r'\[([^\]]+)\]\(([^)]+)\)', r'<a href="\2">\1</a>', text)
        
        return text
    
    def _atomic_save(self, prs: Any, output_path: str) -> None:
        """
        Atomic save presentation
        
        Use temp file + shutil.move to ensure atomic write.
        
        Args:
            prs: Presentation object
            output_path: Final output path
        """
        import shutil
        
        # Create temp file (without file descriptor)
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)
        
        temp_path = tempfile.mktemp(suffix='.pptx', dir=output_dir or '.')
        try:
            prs.save(temp_path)
            # Use shutil.move instead of os.replace, more reliable
            shutil.move(temp_path, output_path)
        except Exception:
            # Clean up temp file
            try:
                os.unlink(temp_path)
            except OSError:
                pass
            raise
    
    def _should_use_template_renderer(self):
        return os.environ.get("USE_TEMPLATE_RENDERER", "0") == "1"

    def _auto_generate_charts(self, slide_data: Dict, template_name: str, chart_gen) -> None:
        """Auto-generate chart images for slides with data but no images."""
        from src.services.chart_generator import ChartConfig, ChartType
        title = slide_data.get("title", "")
        items = slide_data.get("items", [])
        table_data = slide_data.get("table_data", [])
        content = slide_data.get("content", "")
        text_parts = []
        if title:
            text_parts.append(title)
        for item in items:
            text_parts.append(item)
        text_for_analysis = "\n".join(text_parts)
        
        if table_data and not text_for_analysis.strip():
            table_text_parts = []
            for row in table_data:
                row_str = " | ".join(str(c) for c in row)
                table_text_parts.append(row_str)
            text_for_analysis = "\n".join(table_text_parts)
        
        if not text_for_analysis.strip():
            return
        try:
            suggestions = chart_gen.analyze_content(
                section_title=title,
                content=text_for_analysis,
                data_points=None,
            )
        except Exception:
            logger.warning(f"Chart analysis failed for slide: {title}")
            return
        if not suggestions and table_data and len(table_data) >= 2:
            try:
                chart_path = self._generate_chart_from_table(table_data, title, chart_gen)
                if chart_path and os.path.isfile(chart_path):
                    images = slide_data.get("images", [])
                    images.append({"src": chart_path, "alt": title})
                    slide_data["images"] = images
                    logger.info(f"Auto-generated table chart for '{title}': {chart_path}")
            except Exception:
                logger.warning(f"Table chart generation failed for slide: {title}")
            return
        if not suggestions:
            return
        max_charts = 1
        if template_name == "chart_split":
            max_charts = 2
        images = slide_data.get("images", [])
        for suggestion in suggestions[:max_charts]:
            try:
                chart_path = chart_gen.generate_chart(suggestion)
                if chart_path and os.path.isfile(chart_path):
                    images.append({"src": chart_path, "alt": suggestion.title})
                    logger.info(f"Auto-generated chart for '{title}': {chart_path}")
            except Exception:
                logger.warning(f"Chart generation failed for slide: {title}")
        if images:
            slide_data["images"] = images

    def _generate_chart_from_table(self, table_data, title, chart_gen) -> Optional[str]:
        """Generate a bar chart directly from table_data when SmartChartGenerator fails."""
        from src.services.chart_generator import ChartConfig, ChartType
        if len(table_data) < 2:
            return None
        headers = [str(h) for h in table_data[0]]
        categories = []
        values = []
        numeric_col_idx = None
        for col_idx in range(len(headers)):
            numeric_vals = []
            for row_idx in range(1, len(table_data)):
                cell = str(table_data[row_idx][col_idx]) if col_idx < len(table_data[row_idx]) else ""
                cleaned = re.sub(r'[^\d.\-]', '', cell)
                try:
                    numeric_vals.append(float(cleaned))
                except (ValueError, TypeError):
                    numeric_vals.append(None)
            non_none = [v for v in numeric_vals if v is not None]
            if len(non_none) >= len(table_data) * 0.5 and numeric_col_idx is None:
                numeric_col_idx = col_idx
        
        if numeric_col_idx is None:
            return None
        
        name_col = 0 if numeric_col_idx != 0 else (1 if len(headers) > 1 else None)
        if name_col is None:
            return None
        
        for row_idx in range(1, len(table_data)):
            name = str(table_data[row_idx][name_col]) if name_col < len(table_data[row_idx]) else ""
            cell = str(table_data[row_idx][numeric_col_idx]) if numeric_col_idx < len(table_data[row_idx]) else ""
            cleaned = re.sub(r'[^\d.\-]', '', cell)
            try:
                val = float(cleaned)
                categories.append(name[:12])
                values.append(val)
            except (ValueError, TypeError):
                continue
        
        if len(categories) < 2:
            return None
        
        config = ChartConfig(
            chart_type=ChartType.BAR,
            title=title,
            data={"categories": categories, "values": values},
            ylabel=headers[numeric_col_idx],
            source="",
            width=8,
            height=5,
        )
        result = chart_gen.chart_generator.generate(config)
        if result.success and result.image_path:
            dst = chart_gen.output_dir / Path(result.image_path).name
            if Path(result.image_path) != dst and Path(result.image_path).exists():
                import shutil
                shutil.move(str(result.image_path), str(dst))
                return str(dst)
            return result.image_path
        return None

    def _create_pptx_document(
        self,
        slides: List[Dict[str, Any]],
        output_path: str,
        styles: Dict[str, Any]
    ) -> ConversionResult:
        """
        Create PPT document using python-pptx
        
        Args:
            slides: Parsed slide list
            output_path: Output path
            styles: Style configuration
            
        Returns:
            ConversionResult
        """
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions
        prs.slide_width = Inches(styles.get("slide_width", 10))
        prs.slide_height = Inches(styles.get("slide_height", 7.5))
        
        slides_count = 0
        
        if self._should_use_template_renderer():
            from .template_selector import TemplateRegistry, TemplateSelector
            from .slide_renderer import SlideRenderer
            
            registry = TemplateRegistry()
            selector = TemplateSelector()
            renderer = SlideRenderer(self.DESIGN)
            section_index = 0
            
            has_image_slot_cache = {}
            for tname in registry.list_templates():
                has_image_slot_cache[tname] = any(
                    s.get("type") == "image" for s in registry.get(tname).get("slots", [])
                )
            
            try:
                from src.services.smart_chart_generator import SmartChartGenerator
                chart_gen = SmartChartGenerator(output_dir=os.path.join(
                    os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))),
                    "output", "charts"
                ))
            except Exception:
                chart_gen = None
            
            for slide_data in slides:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                slides_count += 1
                
                slide_type = slide_data.get("slide_type", "content")
                if slide_type in ("section_title", "section-title"):
                    section_index += 1
                
                template_name = selector.select_and_enhance(slide_data, section_index=section_index)
                try:
                    template = registry.get(template_name)
                except KeyError:
                    template = registry.get("content_text_only")
                
                source = slide_data.get("source_text", "")
                if source:
                    for dec in template.get("decorations", []):
                        if dec.get("type") == "source_text" and not dec.get("text"):
                            dec["text"] = source
                
                if has_image_slot_cache.get(template_name) and not slide_data.get("images") and chart_gen:
                    self._auto_generate_charts(slide_data, template_name, chart_gen)
                
                renderer.render(slide, slide_data, template, styles, page_num=slides_count)
        else:
            for slide_data in slides:
                slide_type = slide_data.get("slide_type", "content")
                
                # Use blank layout
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                slides_count += 1
                
                # Process by type
                if slide_type == "cover":
                    self._create_cover_slide(slide, slide_data, styles)
                elif slide_type == "toc":
                    self._create_toc_slide(slide, slide_data, styles)
                elif slide_type == "findings":
                    self._create_findings_slide(slide, slide_data, styles)
                elif slide_type == "data":
                    self._create_data_slide(slide, slide_data, styles)
                elif slide_type == "end":
                    self._create_end_slide(slide, slide_data, styles)
                else:
                    self._create_content_slide(slide, slide_data, styles)
        
        # Atomic save
        self._atomic_save(prs, output_path)
        
        # Get file size
        file_size = os.path.getsize(output_path)
        
        logger.info(f"PPT created: {output_path}, size={file_size}, slides={slides_count}")
        
        return ConversionResult(
            success=True,
            output_path=output_path,
            file_size=file_size,
            slides_count=slides_count
        )
    
    # Design system colors (matching ppt_default.html CSS)
    DESIGN = {
        "navy": "1A2744",
        "navy_dark": "0F1A2E",
        "navy_light": "2C3E50",
        "gold": "C9A227",
        "gold_light": "D4AF37",
        "white": "FFFFFF",
        "off_white": "F5F5F5",
        "text_dark": "333333",
        "text_mid": "666666",
        "text_light": "999999",
    }

    def _set_slide_bg(self, slide, color_hex: str):
        """Set solid background color for a slide."""
        from pptx.util import Pt
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb(color_hex)

    def _set_gradient_bg(self, slide, color1_hex: str, color2_hex: str):
        """Set gradient background for a slide."""
        bg = slide.background
        fill = bg.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = self._rgb(color1_hex)
        fill.gradient_stops[1].color.rgb = self._rgb(color2_hex)

    def _rgb(self, hex_color: str):
        """Convert hex color string to RGBColor."""
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        hex_color = hex_color.lstrip('#')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

    def _add_footer_bar(self, slide, styles: Dict[str, Any], color: str = None):
        """Add gold footer bar at bottom of slide."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        sw = styles.get("slide_width", 13.333)
        sh = styles.get("slide_height", 7.5)
        bar_color = color or self.DESIGN["gold"]
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(sh - 0.11), Inches(sw), Inches(0.11)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(bar_color)
        shape.line.fill.background()

    def _add_title_underline(self, slide, left: float, top: float, width: float):
        """Add gold underline below section title."""
        from pptx.util import Inches
        shape = slide.shapes.add_shape(
            1, Inches(left), Inches(top), Inches(width), Inches(0.04)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(self.DESIGN["gold"])
        shape.line.fill.background()

    def _add_side_accent(self, slide, styles: Dict[str, Any]):
        """Add vertical gold accent bar on left side of content slides."""
        from pptx.util import Inches
        sh = styles.get("slide_height", 7.5)
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(0.06), Inches(sh)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(self.DESIGN["gold"])
        shape.line.fill.background()

    def _content_width(self, styles: Dict[str, Any], margin: float = 1.0) -> float:
        """Calculate content width in inches based on slide_width minus margins."""
        return styles.get("slide_width", 13.333) - margin

    def _create_cover_slide(self, slide, slide_data: Dict[str, Any], styles: Dict[str, Any]):
        """Create cover slide with deep navy gradient + gold title."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        try:
            self._set_gradient_bg(slide, self.DESIGN["navy"], self.DESIGN["navy_light"])
        except Exception:
            self._set_slide_bg(slide, self.DESIGN["navy"])

        cw = self._content_width(styles)
        sw = styles.get("slide_width", 13.333)

        self._add_footer_bar(slide, styles, self.DESIGN["gold"])

        title = slide_data.get("title", "")
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(sw - 2), Inches(1.5)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(styles.get("title_size", 44))
            p.font.bold = True
            p.font.color.rgb = self._rgb(self.DESIGN["gold"])
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER

        subtitle = slide_data.get("subtitle", "")
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2), Inches(sw - 2), Inches(0.8)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = self._rgb(self.DESIGN["white"])
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER

        import datetime
        date_box = slide.shapes.add_textbox(
            Inches(1), Inches(5.5), Inches(sw - 2), Inches(0.5)
        )
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = datetime.date.today().strftime("%Y-%m-%d")
        p.font.size = Pt(18)
        p.font.color.rgb = self._rgb(self.DESIGN["white"])
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER

    def _create_toc_slide(self, slide, slide_data: Dict[str, Any], styles: Dict[str, Any]):
        """Create TOC slide with gold left accent on items."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        self._set_slide_bg(slide, self.DESIGN["off_white"])
        cw = self._content_width(styles)

        self._add_footer_bar(slide, styles)

        title = slide_data.get("title", "目录")
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.5), Inches(cw), Inches(1)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(styles.get("subtitle_size", 28))
            p.font.bold = True
            p.font.color.rgb = self._rgb(self.DESIGN["navy"])
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER

        items = slide_data.get("items", [])
        if items:
            content_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(2), Inches(cw - 0.5), Inches(4)
            )
            tf = content_box.text_frame
            for i, item in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.space_after = Pt(12)
                run = p.add_run()
                run.text = "  " + item
                run.font.size = Pt(20)
                run.font.color.rgb = self._rgb(self.DESIGN["text_dark"])
                run.font.name = "Microsoft YaHei"
                bullet = p.add_run()
                bullet.text = "■ "
                bullet.font.size = Pt(12)
                bullet.font.color.rgb = self._rgb(self.DESIGN["gold"])
                bullet.font.name = "Microsoft YaHei"

            left_bar = slide.shapes.add_shape(
                1, Inches(1.0), Inches(2.0), Inches(0.05), Inches(min(len(items) * 0.6, 4.0))
            )
            left_bar.fill.solid()
            left_bar.fill.fore_color.rgb = self._rgb(self.DESIGN["gold"])
            left_bar.line.fill.background()
    
    def _add_images_to_slide(self, slide, slide_data: Dict[str, Any], styles: Dict[str, Any], position: str = "right"):
        """Add image shapes to slide, preserving aspect ratio.
        
        position: "right" = right half of slide (left-text-right-image layout)
                  "below" = below text area (top-text-bottom-image layout)
        """
        from pptx.util import Inches, Pt
        
        images = slide_data.get("images", [])
        if not images:
            return
        
        slide_width = styles.get("slide_width", 13.333)
        slide_height = styles.get("slide_height", 7.5)
        
        if position == "right":
            img_area_left = slide_width * 0.48
            img_area_top = 1.6
            img_area_w = slide_width - img_area_left - 0.5
            img_area_h = slide_height - img_area_top - 0.8
            
            n_images = len(images)
            if n_images == 1:
                self._place_image(slide, images[0], img_area_left, img_area_top, img_area_w, img_area_h)
            else:
                per_height = (img_area_h - 0.3 * (n_images - 1)) / n_images
                for idx, img_info in enumerate(images):
                    top = img_area_top + idx * (per_height + 0.3)
                    self._place_image(slide, img_info, img_area_left, top, img_area_w, per_height)
        else:
            max_w = min(9.0, slide_width - 2.0)
            max_h = min(3.5, slide_height - 4.5)
            top = 4.0
            left = (slide_width - max_w) / 2 + 0.3
            n_images = len(images)
            if n_images == 1:
                self._place_image(slide, images[0], left, top, max_w, max_h)
            else:
                cols = 2
                per_w = (max_w - 0.5) / cols
                for idx, img_info in enumerate(images):
                    r = idx // cols
                    c = idx % cols
                    il = left + c * (per_w + 0.5)
                    it = top + r * (max_h + 0.3)
                    self._place_image(slide, img_info, il, it, per_w, max_h)
    
    def _place_image(self, slide, img_info: Dict[str, str], left: float, top: float, max_w: float, max_h: float):
        """Place a single image preserving aspect ratio within max_w x max_h bounds."""
        from pptx.util import Inches
        
        src = img_info.get("src", "")
        if not src:
            return
        
        if not os.path.isfile(src):
            logger.warning(f"Image file not found, skipping: {src}")
            return
        
        try:
            from PIL import Image as PILImage
            with PILImage.open(src) as pil_img:
                img_w, img_h = pil_img.size
                aspect = img_w / img_h
            
            if max_w / max_h > aspect:
                final_h = max_h
                final_w = final_h * aspect
            else:
                final_w = max_w
                final_h = final_w / aspect
            
            slide.shapes.add_picture(src, Inches(left), Inches(top), Inches(final_w), Inches(final_h))
        except ImportError:
            try:
                slide.shapes.add_picture(src, Inches(left), Inches(top), Inches(max_w))
            except Exception as e:
                logger.warning(f"Failed to add image to slide: {src}, error: {e}")
        except Exception as e:
            logger.warning(f"Failed to add image to slide: {src}, error: {e}")

    def _create_content_slide(self, slide, slide_data: Dict[str, Any], styles: Dict[str, Any]):
        """Create content slide: title bar + left text / right image layout."""
        from pptx.util import Inches, Pt

        self._set_slide_bg(slide, self.DESIGN["white"])
        sw = styles.get("slide_width", 13.333)
        sh = styles.get("slide_height", 7.5)
        self._add_side_accent(slide, styles)
        self._add_footer_bar(slide, styles)

        title = slide_data.get("title", "")
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.3), Inches(sw - 1.6), Inches(0.7)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self._rgb(self.DESIGN["navy"])
            p.font.name = "Microsoft YaHei"
            self._add_title_underline(slide, 0.8, 1.05, min(4.0, sw * 0.3))

        images = slide_data.get("images", [])
        has_image = bool(images)

        if has_image:
            text_left = 0.8
            text_top = 1.3
            text_width = sw * 0.42
            text_height = sh - text_top - 1.0
        else:
            text_left = 0.8
            text_top = 1.3
            text_width = sw - 1.6
            text_height = sh - text_top - 1.0

        content = slide_data.get("content", "")
        items = slide_data.get("items", [])

        if items:
            content_box = slide.shapes.add_textbox(
                Inches(text_left), Inches(text_top), Inches(text_width), Inches(text_height)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            for i, item in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                run_bullet = p.add_run()
                run_bullet.text = "▸ "
                run_bullet.font.size = Pt(14)
                run_bullet.font.color.rgb = self._rgb(self.DESIGN["gold"])
                run_bullet.font.name = "Microsoft YaHei"
                run_text = p.add_run()
                run_text.text = item
                run_text.font.size = Pt(14)
                run_text.font.color.rgb = self._rgb(self.DESIGN["text_dark"])
                run_text.font.name = "Microsoft YaHei"
                p.space_after = Pt(6)
                p.space_before = Pt(2)
        elif content:
            content_box = slide.shapes.add_textbox(
                Inches(text_left), Inches(text_top), Inches(text_width), Inches(text_height)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(14)
            p.font.color.rgb = self._rgb(self.DESIGN["text_dark"])
            p.font.name = "Microsoft YaHei"
            p.space_after = Pt(6)

        if has_image:
            self._add_images_to_slide(slide, slide_data, styles, position="right")

    def _create_findings_slide(self, slide, slide_data: Dict[str, Any], styles: Dict[str, Any]):
        """Create findings slide: left check-mark items / right image."""
        from pptx.util import Inches, Pt

        self._set_slide_bg(slide, self.DESIGN["white"])
        sw = styles.get("slide_width", 13.333)
        sh = styles.get("slide_height", 7.5)
        self._add_side_accent(slide, styles)
        self._add_footer_bar(slide, styles)

        title = slide_data.get("title", "")
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.3), Inches(sw - 1.6), Inches(0.7)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self._rgb(self.DESIGN["navy"])
            p.font.name = "Microsoft YaHei"
            self._add_title_underline(slide, 0.8, 1.05, min(4.0, sw * 0.3))

        images = slide_data.get("images", [])
        has_image = bool(images)

        if has_image:
            text_width = sw * 0.42
        else:
            text_width = sw - 1.6

        items = slide_data.get("items", [])
        if items:
            content_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.3), Inches(text_width), Inches(sh - 2.3)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            for i, item in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                run_check = p.add_run()
                run_check.text = "✓ "
                run_check.font.size = Pt(16)
                run_check.font.color.rgb = self._rgb(self.DESIGN["gold"])
                run_check.font.bold = True
                run_check.font.name = "Microsoft YaHei"
                run_text = p.add_run()
                run_text.text = item
                run_text.font.size = Pt(14)
                run_text.font.color.rgb = self._rgb(self.DESIGN["text_dark"])
                run_text.font.name = "Microsoft YaHei"
                p.space_after = Pt(6)

        if has_image:
            self._add_images_to_slide(slide, slide_data, styles, position="right")

    def _create_data_slide(self, slide, slide_data: Dict[str, Any], styles: Dict[str, Any]):
        """Create data slide: left table / right chart."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        self._set_slide_bg(slide, self.DESIGN["white"])
        sw = styles.get("slide_width", 13.333)
        sh = styles.get("slide_height", 7.5)
        self._add_side_accent(slide, styles)
        self._add_footer_bar(slide, styles)

        title = slide_data.get("title", "")
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.3), Inches(sw - 1.6), Inches(0.7)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = self._rgb(self.DESIGN["navy"])
            p.font.name = "Microsoft YaHei"
            self._add_title_underline(slide, 0.8, 1.05, min(4.0, sw * 0.3))

        images = slide_data.get("images", [])
        has_image = bool(images)

        if has_image:
            table_width = sw * 0.42
        else:
            table_width = sw - 1.6

        table_data = slide_data.get("table_data", [])
        if table_data:
            rows = len(table_data)
            cols = max(len(row) for row in table_data) if table_data else 0
            if rows > 0 and cols > 0:
                table = slide.shapes.add_table(
                    rows, cols, Inches(0.8), Inches(1.3), Inches(table_width), Inches(0.4 * rows)
                ).table
                for i, row_data in enumerate(table_data):
                    for j, cell_text in enumerate(row_data):
                        if j < cols:
                            import re
                            clean = re.sub(r'</?(strong|em|code|del|a[^>]*)>', '', str(cell_text))
                            cell = table.cell(i, j)
                            cell.text = clean
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.font.size = Pt(12)
                                paragraph.font.name = "Microsoft YaHei"
                                if i == 0:
                                    paragraph.font.color.rgb = self._rgb(self.DESIGN["white"])
                                    paragraph.font.bold = True
                                else:
                                    paragraph.font.color.rgb = self._rgb(self.DESIGN["text_dark"])
                            if i == 0:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = self._rgb(self.DESIGN["navy"])
                            elif i % 2 == 0:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = self._rgb(self.DESIGN["off_white"])

        if has_image:
            self._add_images_to_slide(slide, slide_data, styles, position="right")

    def _create_end_slide(self, slide, slide_data: Dict[str, Any], styles: Dict[str, Any]):
        """Create end slide with navy background + gold title."""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        self._set_slide_bg(slide, self.DESIGN["navy"])
        cw = self._content_width(styles)
        sw = styles.get("slide_width", 13.333)

        self._add_footer_bar(slide, styles, self.DESIGN["white"])

        title = slide_data.get("title", "")
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(sw - 2), Inches(1.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = self._rgb(self.DESIGN["gold"])
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER

        thanks_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.5), Inches(sw - 2), Inches(0.8)
        )
        tf = thanks_box.text_frame
        p = tf.paragraphs[0]
        p.text = "感谢关注"
        p.font.size = Pt(28)
        p.font.color.rgb = self._rgb(self.DESIGN["white"])
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER
    
    def _add_formatted_text(self, paragraph, text: str, styles: Dict[str, Any]) -> None:
        """
        Add formatted text to paragraph
        
        Parse HTML inline tags in text (converted from Markdown)
        <strong>bold</strong> → bold
        <em>italic</em> → italic
        <code>code</code> → code
        
        Args:
            paragraph: Paragraph object
            text: Text that may contain HTML tags
            styles: Style configuration
        """
        import re
        from pptx.util import Pt
        
        # Clean HTML tags, extract plain text and format info
        # Simplified: directly clean HTML tags, PPT does not support complex inline formatting
        clean_text = re.sub(r'</?(strong|em|code|del|a[^>]*)>', '', text)
        
        # Check for bold markers
        has_bold = '<strong>' in text or '**' in text
        
        # Check for italic markers
        has_italic = ('<em>' in text or '*' in text) and '**' not in text
        
        # Check for code markers
        has_code = '<code>' in text or '`' in text
        
        paragraph.text = clean_text
        
        # Apply format (entire paragraph uses same format in PPT)
        if has_bold:
            paragraph.font.bold = True
        if has_italic:
            paragraph.font.italic = True
        if has_code:
            paragraph.font.name = 'Courier New'
            paragraph.font.size = Pt(styles.get("body_size", 18) - 2)
    
    def _create_fallback_document(
        self,
        slides: List[Dict[str, Any]],
        output_path: str,
        styles: Dict[str, Any]
    ) -> ConversionResult:
        """
        Create fallback format document (when python-pptx is unavailable)
        """
        # Create plain text version
        text_lines = []
        
        for i, slide in enumerate(slides, 1):
            slide_type = slide.get("slide_type", "content")
            text_lines.append(f"\n=== Slide {i} ({slide_type}) ===")
            
            if slide.get("title"):
                text_lines.append(f"Title: {slide['title']}")
            
            if slide.get("content"):
                text_lines.append(f"Content: {slide['content']}")
            
            items = slide.get("items", [])
            if items:
                text_lines.append("List items:")
                for item in items:
                    text_lines.append(f"  - {item}")
            
            table_data = slide.get("table_data", [])
            if table_data:
                text_lines.append("Table:")
                for row in table_data:
                    text_lines.append(f"  | {' | '.join(row)} |")
        
        content = "\n".join(text_lines)
        
        # Change to .txt extension
        txt_path = output_path.rsplit('.', 1)[0] + '.txt'
        
        # Atomic write
        fd, temp_path = tempfile.mkstemp(suffix='.txt')
        try:
            with os.fdopen(fd, 'w', encoding='utf-8') as f:
                f.write(content)
            os.replace(temp_path, txt_path)
        except Exception:
            try:
                os.unlink(temp_path)
            except OSError:
                pass
            raise
        
        file_size = os.path.getsize(txt_path)
        
        logger.warning(f"python-pptx not available, created text file: {txt_path}")
        
        return ConversionResult(
            success=True,
            output_path=txt_path,
            file_size=file_size,
            slides_count=len(slides),
            error="python-pptx not available, created text file instead"
        )


# Export
__all__ = ["HTMLToPPTConverter", "ConversionResult"]