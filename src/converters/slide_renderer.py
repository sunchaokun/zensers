import logging
import os
import re
import datetime
from typing import Dict, List, Optional
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)


class SlideRenderer:
    def __init__(self, design: Dict[str, str], image_provider=None):
        self.design = design
        self.image_provider = image_provider

    def render(self, slide, slide_data: Dict, template: Dict, styles: Dict, page_num: int = 0, layout_overrides: Dict = None):
        decorations = template.get("decorations", [])

        self._render_background(slide, template.get("background", {}))

        for dec in decorations:
            if dec.get("layer") != "top":
                self._render_decoration(slide, dec, styles, page_num)

        chart_height = None
        if layout_overrides:
            for content_key in ("data_table", "kpi_row", "bullet_items"):
                if content_key in layout_overrides:
                    chart_height = layout_overrides[content_key].get("height")
                    break

        for slot in template.get("slots", []):
            merged_slot = self._apply_layout_override(slot, layout_overrides)
            if chart_height and slot.get("type") == "image" and slot.get("id") == "chart":
                if "chart" not in (layout_overrides or {}):
                    merged_slot = dict(merged_slot)
                    merged_pos = dict(merged_slot.get("position", {}))
                    merged_pos["height"] = chart_height
                    merged_slot["position"] = merged_pos
            self._render_slot(slide, merged_slot, slide_data, styles)

        for dec in decorations:
            if dec.get("layer") == "top":
                self._render_decoration(slide, dec, styles, page_num)

    def _apply_layout_override(self, slot: Dict, layout_overrides: Dict = None) -> Dict:
        if not layout_overrides:
            return slot
        slot_id = slot.get("id", "")
        override = layout_overrides.get(slot_id)
        if not override:
            return slot
        merged = dict(slot)
        merged_pos = dict(slot.get("position", {}))
        override_copy = dict(override)
        style_delta = override_copy.pop("_style_delta", None)
        merged_pos.update(override_copy)
        merged["position"] = merged_pos
        if style_delta:
            merged_style = dict(slot.get("style", {}))
            merged_style.update(style_delta)
            merged["style"] = merged_style
        return merged

    def _render_slot(self, slide, slot: Dict, slide_data: Dict, styles: Dict):
        slot_type = slot.get("type", "text")
        dispatch = {
            "text": self._render_text_slot,
            "items": self._render_items_slot,
            "image": self._render_image_slot,
            "table": self._render_table_slot,
            "kpi_cards": self._render_kpi_cards_slot,
            "insight_bar": self._render_insight_bar_slot,
            "comparison": self._render_comparison_slot,
        }
        handler = dispatch.get(slot_type)
        if handler:
            handler(slide, slot, slide_data, styles)

    def _resolve_color(self, color_str: str) -> str:
        if color_str in self.design:
            return self.design[color_str]
        return color_str

    def _rgb(self, hex_color: str):
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 3:
            hex_color = hex_color[0]*2 + hex_color[1]*2 + hex_color[2]*2
        if len(hex_color) < 6:
            hex_color = hex_color.ljust(6, '0')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

    def _clean_html(self, text: str) -> str:
        return re.sub(r'</?(strong|em|code|del|a[^>]*)>', '', str(text))

    def _render_background(self, slide, bg: Dict):
        bg_type = bg.get("type", "")
        if bg_type == "solid":
            color = self._resolve_color(bg.get("color", "white"))
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = self._rgb(color)
        elif bg_type == "gradient":
            c1 = self._resolve_color(bg.get("color1", "white"))
            c2 = self._resolve_color(bg.get("color2", "white"))
            fill = slide.background.fill
            fill.gradient()
            fill.gradient_stops[0].position = 0.0
            fill.gradient_stops[0].color.rgb = self._rgb(c1)
            fill.gradient_stops[1].position = 1.0
            fill.gradient_stops[1].color.rgb = self._rgb(c2)

    def _render_text_slot(self, slide, slot: Dict, slide_data: Dict, styles: Dict):
        source = slot.get("source", "title")
        text = ""
        if source == "section_number":
            num = slide_data.get("section_number")
            fmt = slot.get("style", {}).get("format", "")
            if num is not None:
                if fmt == "zero_padded" and isinstance(num, int):
                    text = f"{num:02d}"
                else:
                    text = str(num)
        elif source == "section_summary":
            text = slide_data.get("section_summary", "")
        elif source == "auto":
            auto_type = slot.get("style", {}).get("auto_type", "date")
            if auto_type == "date":
                text = datetime.date.today().strftime("%Y-%m-%d")
            else:
                text = ""
        else:
            text = slide_data.get(source, "")
        if not text:
            return
        pos = slot.get("position", {})
        st = slot.get("style", {})
        left = Inches(pos.get("left", 0.8))
        top = Inches(pos.get("top", 0.3))
        width = Inches(pos.get("width", 11.7))
        height = Inches(pos.get("height", 0.7))
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        lines = str(text).split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(st.get("font_size", 24))
            p.font.bold = st.get("font_weight", "") == "bold"
            color = self._resolve_color(st.get("color", "text_dark"))
            p.font.color.rgb = self._rgb(color)
            p.font.name = st.get("font", "Microsoft YaHei")
            alignment = st.get("alignment", "left")
            align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
            p.alignment = align_map.get(alignment, PP_ALIGN.LEFT)

    def _render_items_slot(self, slide, slot: Dict, slide_data: Dict, styles: Dict):
        items = slide_data.get(slot.get("source", "items"), [])
        if not items:
            return
        pos = slot.get("position", {})
        st = slot.get("style", {})
        max_items = st.get("max_items", 7)
        items = items[:max_items]
        left = Inches(pos.get("left", 0.8))
        top = Inches(pos.get("top", 1.3))
        width = Inches(pos.get("width", 11.7))
        height = Inches(pos.get("height", 5.2))
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        bullet = st.get("bullet", "▸")
        bullet_color = self._resolve_color(st.get("bullet_color", "gold"))
        text_color = self._resolve_color(st.get("color", "text_dark"))
        font_name = st.get("font", "Microsoft YaHei")
        font_size = st.get("font_size", 18)
        line_spacing = st.get("line_spacing", 14)
        bullet_font_size = st.get("bullet_font_size", font_size + 4)
        space_before = st.get("space_before", line_spacing // 3)
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run_bullet = p.add_run()
            run_bullet.text = f"{bullet}  "
            run_bullet.font.size = Pt(bullet_font_size)
            run_bullet.font.color.rgb = self._rgb(bullet_color)
            run_bullet.font.name = font_name
            run_text = p.add_run()
            run_text.text = self._clean_html(item)
            run_text.font.size = Pt(font_size)
            run_text.font.color.rgb = self._rgb(text_color)
            run_text.font.name = font_name
            p.space_after = Pt(line_spacing)
            p.space_before = Pt(space_before)

    def _render_image_slot(self, slide, slot: Dict, slide_data: Dict, styles: Dict):
        images = slide_data.get(slot.get("source", "images"), [])
        idx = slot.get("style", {}).get("index", 0)
        slot_id = slot.get("id", "")
        if slot_id == "chart" and images:
            chart_idx = None
            for ci, img in enumerate(images):
                src = img.get("src", "")
                if img.get("image_type") == "chart" or any(p in src for p in ("pie_", "bar_", "hbar_", "line_", "chart")):
                    chart_idx = ci
                    break
            if chart_idx is not None:
                idx = chart_idx
        if idx >= len(images):
            return
        img_info = images[idx]
        src = img_info.get("src", "")
        if not src:
            return
        if not os.path.isfile(src):
            if self.image_provider and src.startswith(("http://", "https://")):
                resolved = self.image_provider.resolve_image_src(src)
                if resolved:
                    src = resolved
                else:
                    logger.warning(f"Failed to resolve image URL, skipping: {src[:80]}")
                    return
            else:
                logger.warning(f"Image file not found, skipping: {src}")
                return
        pos = slot.get("position", {})
        max_w = pos.get("width", 6.4)
        max_h = pos.get("height", 5.2)
        slot_left = pos.get("left", 0.8)
        slot_top = pos.get("top", 1.6)
        is_chart = img_info.get("image_type") == "chart" or "chart" in src or "pie_" in src or "bar_" in src or "hbar_" in src or "line_" in src or slot_id == "chart"
        if is_chart:
            slide.shapes.add_picture(src, Inches(slot_left), Inches(slot_top), Inches(max_w), Inches(max_h))
        else:
            try:
                from PIL import Image as PILImage
                with PILImage.open(src) as pil_img:
                    img_w, img_h = pil_img.size
                    aspect = img_w / img_h
                if max_w / max_h > aspect:
                    final_h = max_h
                    final_w = final_h * aspect
                else:
                    final_w = max_w
                    final_h = final_w / aspect
                centered_left = slot_left + (max_w - final_w) / 2
                centered_top = slot_top + (max_h - final_h) / 2
                slide.shapes.add_picture(src, Inches(centered_left), Inches(centered_top), Inches(final_w), Inches(final_h))
            except ImportError:
                try:
                    slide.shapes.add_picture(src, Inches(slot_left), Inches(slot_top), Inches(max_w))
                except Exception as e:
                    logger.warning(f"Failed to add image '{src}': {e}")
            except Exception as e:
                logger.warning(f"Failed to add image '{src}': {e}")

    def _render_table_slot(self, slide, slot: Dict, slide_data: Dict, styles: Dict):
        table_data = slide_data.get(slot.get("source", "table_data"), [])
        if not table_data:
            return
        pos = slot.get("position", {})
        st = slot.get("style", {})
        has_images = bool(slide_data.get("images", []))
        width_val = pos.get("width", 5.6)
        if width_val == "auto":
            slide_w = styles.get("slide_width", 13.33)
            width_val = slide_w - pos.get("left", 0.8) - 0.8
        elif not has_images and st.get("full_width_when_no_image", False):
            slide_w = styles.get("slide_width", 13.33)
            width_val = slide_w - pos.get("left", 0.8) - 0.8
        height = pos.get("height", "auto")
        row_height = st.get("row_height", 0.55)
        if height == "auto":
            height = len(table_data) * row_height
        else:
            height = float(height)
        left = Inches(pos.get("left", 0.8))
        top = Inches(pos.get("top", 1.3))
        width = Inches(width_val)
        rows = len(table_data)
        cols = max(len(r) for r in table_data) if table_data else 1
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(height))
        table = table_shape.table
        header_bg = self._resolve_color(st.get("header_bg", "navy"))
        header_color = self._resolve_color(st.get("header_color", "white"))
        stripe = st.get("stripe", True)
        stripe_color = self._resolve_color(st.get("stripe_color", "off_white"))
        font_name = st.get("font", "Microsoft YaHei")
        row_font_size = st.get("row_font_size", 14)
        header_font_size = st.get("header_font_size", row_font_size + 2)
        cell_padding = st.get("cell_padding", 4)
        for c_idx in range(cols):
            table.columns[c_idx].width = Inches(width_val / cols)
        for r_idx, row_data in enumerate(table_data):
            for c_idx in range(cols):
                cell = table.cell(r_idx, c_idx)
                raw_text = row_data[c_idx] if c_idx < len(row_data) else ""
                cell.text = self._clean_html(str(raw_text))
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                for paragraph in cell.text_frame.paragraphs:
                    if r_idx == 0:
                        paragraph.font.size = Pt(header_font_size)
                        paragraph.font.color.rgb = self._rgb(header_color)
                        paragraph.font.bold = True
                        paragraph.alignment = PP_ALIGN.CENTER
                    else:
                        paragraph.font.size = Pt(row_font_size)
                        paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.font.name = font_name
                    paragraph.space_before = Pt(cell_padding)
                    paragraph.space_after = Pt(cell_padding)
                if r_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self._rgb(header_bg)
                elif stripe and r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self._rgb(stripe_color)

    def _render_kpi_cards_slot(self, slide, slot: Dict, slide_data: Dict, styles: Dict):
        kpi_data = slide_data.get(slot.get("source", "kpi_data"), [])
        if not kpi_data:
            return
        pos = slot.get("position", {})
        st = slot.get("style", {})
        max_cards = st.get("max_cards", 4)
        kpi_data = kpi_data[:max_cards]
        if not kpi_data:
            return
        layout_mode = st.get("layout_mode", "row")
        grid_cols = st.get("grid_cols", 2)
        total_width = pos.get("width", 11.7)
        gap = st.get("card_gap", 0.4)
        card_height = pos.get("height", 3.5)
        start_left = pos.get("left", 0.8)
        start_top = pos.get("top", 1.5)
        card_bg = self._resolve_color(st.get("card_bg", "navy"))
        number_size = st.get("number_size", 36)
        number_color = self._resolve_color(st.get("number_color", "gold"))
        label_size = st.get("label_size", 12)
        label_color = self._resolve_color(st.get("label_color", "white"))
        if layout_mode == "grid":
            cols = min(grid_cols, len(kpi_data))
            card_width = (total_width - gap * (cols - 1)) / cols
            rows = (len(kpi_data) + cols - 1) // cols
            grid_card_height = (card_height - gap * (rows - 1)) / rows if rows > 1 else card_height
            for i, kpi in enumerate(kpi_data):
                r = i // cols
                c = i % cols
                card_left = start_left + c * (card_width + gap)
                card_top = start_top + r * (grid_card_height + gap)
                self._render_single_kpi_card(slide, kpi, card_left, card_top, card_width, grid_card_height, card_bg, number_size, number_color, label_size, label_color, st)
        else:
            card_width = (total_width - gap * (len(kpi_data) - 1)) / len(kpi_data)
            for i, kpi in enumerate(kpi_data):
                card_left = start_left + i * (card_width + gap)
                self._render_single_kpi_card(slide, kpi, card_left, start_top, card_width, card_height, card_bg, number_size, number_color, label_size, label_color, st)

    def _render_single_kpi_card(self, slide, kpi, card_left, card_top, card_width, card_height, card_bg, number_size, number_color, label_size, label_color, st):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(card_left), Inches(card_top),
            Inches(card_width), Inches(card_height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(card_bg)
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        number_text = kpi.get("number", "")
        trend = kpi.get("trend")
        trend_dir = kpi.get("trend_direction")
        if trend_dir and trend:
            arrow = st.get("trend_up", "↑") if trend_dir == "up" else st.get("trend_down", "↓")
            trend_color_key = "trend_color_up" if trend_dir == "up" else "trend_color_down"
            trend_color = self._resolve_color(st.get(trend_color_key, "4CAF50"))
            p_num = tf.paragraphs[0]
            run_num = p_num.add_run()
            run_num.text = number_text + " "
            run_num.font.size = Pt(number_size)
            run_num.font.color.rgb = self._rgb(number_color)
            run_num.font.bold = True
            run_num.font.name = "Microsoft YaHei"
            run_trend = p_num.add_run()
            run_trend.text = f"{arrow} {trend}"
            run_trend.font.size = Pt(16)
            run_trend.font.color.rgb = self._rgb(trend_color)
            run_trend.font.name = "Microsoft YaHei"
        else:
            p_num = tf.paragraphs[0]
            p_num.text = number_text
            p_num.font.size = Pt(number_size)
            p_num.font.color.rgb = self._rgb(number_color)
            p_num.font.bold = True
            p_num.font.name = "Microsoft YaHei"
        label = kpi.get("label", "")
        if not label:
            label = kpi.get("original_text", "")[:30]
        if label:
            p_label = tf.add_paragraph()
            p_label.text = label
            p_label.font.size = Pt(label_size)
            p_label.font.color.rgb = self._rgb(label_color)
            p_label.font.name = "Microsoft YaHei"
            p_label.alignment = PP_ALIGN.CENTER

    def _render_insight_bar_slot(self, slide, slot: Dict, slide_data: Dict, styles: Dict):
        text = slide_data.get(slot.get("source", "insight_text"), "")
        if not text:
            return
        pos = slot.get("position", {})
        st = slot.get("style", {})
        left = Inches(pos.get("left", 0.8))
        top = Inches(pos.get("top", 5.8))
        width = Inches(pos.get("width", 11.7))
        height = Inches(pos.get("height", 1.0))
        bg_color = self._resolve_color(st.get("bg_color", "navy"))
        icon = st.get("icon", "")
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(bg_color)
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        font_name = st.get("font", "Microsoft YaHei")
        font_color = self._resolve_color(st.get("color", "white"))
        if icon:
            run_icon = p.add_run()
            run_icon.text = f"{icon}  "
            run_icon.font.size = Pt(st.get("font_size", 13))
            run_icon.font.color.rgb = self._rgb(font_color)
            run_icon.font.name = font_name
        run_text = p.add_run()
        run_text.text = text
        run_text.font.size = Pt(st.get("font_size", 13))
        run_text.font.color.rgb = self._rgb(font_color)
        run_text.font.name = font_name

    def _render_comparison_slot(self, slide, slot: Dict, slide_data: Dict, styles: Dict):
        comp_data = slide_data.get(slot.get("source", "comparison_data"), {})
        if not comp_data:
            return
        pos = slot.get("position", {})
        st = slot.get("style", {})
        left_side = comp_data.get("left", {})
        right_side = comp_data.get("right", {})
        total_width = pos.get("width", 11.7)
        col_width = (total_width - 0.4) / 2
        start_left = pos.get("left", 0.8)
        start_top = pos.get("top", 1.3)
        height = pos.get("height", 5.2)
        available = height - 0.6
        max_items = int(available / 0.4)
        left_color = self._resolve_color(st.get("left_color", "navy"))
        right_color = self._resolve_color(st.get("right_color", "gold"))
        font_size = st.get("font_size", 13)
        font_name = st.get("font", "Microsoft YaHei")
        for col_idx, (side, color) in enumerate([(left_side, left_color), (right_side, right_color)]):
            col_left = start_left + col_idx * (col_width + 0.4)
            title = side.get("title", "")
            items = side.get("items", [])[:max_items]
            y = start_top
            if title:
                title_shape = slide.shapes.add_textbox(
                    Inches(col_left), Inches(y), Inches(col_width), Inches(0.5)
                )
                tf = title_shape.text_frame
                tf.paragraphs[0].text = title
                tf.paragraphs[0].font.size = Pt(18)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = self._rgb(color)
                tf.paragraphs[0].font.name = font_name
                y += 0.6
            for item in items:
                item_shape = slide.shapes.add_textbox(
                    Inches(col_left), Inches(y), Inches(col_width), Inches(0.4)
                )
                tf = item_shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                run_b = p.add_run()
                run_b.text = "▸  "
                run_b.font.size = Pt(font_size)
                run_b.font.color.rgb = self._rgb(self._resolve_color("gold"))
                run_b.font.name = font_name
                run_t = p.add_run()
                run_t.text = self._clean_html(item)
                run_t.font.size = Pt(font_size)
                run_t.font.color.rgb = self._rgb(self._resolve_color("text_dark"))
                run_t.font.name = font_name
                y += 0.4

    def _render_decoration(self, slide, dec: Dict, styles: Dict, page_num: int = 0):
        dec_type = dec.get("type", "")
        font_name = "Microsoft YaHei"

        if dec_type == "footer_bar":
            color = self._resolve_color(dec.get("color", "gold"))
            height = dec.get("height", 0.11)
            slide_h = styles.get("slide_height", 7.5)
            slide_w = styles.get("slide_width", 13.33)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(slide_h - height),
                Inches(slide_w), Inches(height),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(color)
            shape.line.fill.background()

        elif dec_type == "side_accent":
            color = self._resolve_color(dec.get("color", "gold"))
            width = dec.get("width", 0.06)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(width), Inches(styles.get("slide_height", 7.5)),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(color)
            shape.line.fill.background()

        elif dec_type == "left_accent":
            color = self._resolve_color(dec.get("color", "gold"))
            w = dec.get("width", 0.05)
            h = dec.get("height", 4.0)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.3), Inches(1.5),
                Inches(w), Inches(h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(color)
            shape.line.fill.background()

        elif dec_type == "title_underline":
            color = self._resolve_color(dec.get("color", "gold"))
            w = dec.get("width", 4.0)
            offset_top = dec.get("offset_top", 1.05)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.8), Inches(offset_top),
                Inches(w), Inches(0.04),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(color)
            shape.line.fill.background()

        elif dec_type == "page_number":
            if page_num > 0:
                position = dec.get("position", "bottom_right")
                color = self._resolve_color(dec.get("color", "text_light"))
                font_size = dec.get("font_size", 10)
                slide_w = styles.get("slide_width", 13.33)
                slide_h = styles.get("slide_height", 7.5)
                if position == "bottom_right":
                    left = Inches(slide_w - 1.5)
                    top = Inches(slide_h - 0.5)
                else:
                    left = Inches(0.5)
                    top = Inches(slide_h - 0.5)
                txBox = slide.shapes.add_textbox(left, top, Inches(1.2), Inches(0.3))
                tf = txBox.text_frame
                tf.paragraphs[0].text = str(page_num)
                tf.paragraphs[0].font.size = Pt(font_size)
                tf.paragraphs[0].font.color.rgb = self._rgb(color)
                tf.paragraphs[0].font.name = font_name
                tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

        elif dec_type == "source_text":
            text = dec.get("text", "")
            if not text:
                return
            color = self._resolve_color(dec.get("color", "text_light"))
            font_size = dec.get("font_size", 9)
            slide_h = styles.get("slide_height", 7.5)
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(slide_h - 0.5), Inches(4.0), Inches(0.3))
            tf = txBox.text_frame
            tf.paragraphs[0].text = text
            tf.paragraphs[0].font.size = Pt(font_size)
            tf.paragraphs[0].font.color.rgb = self._rgb(color)
            tf.paragraphs[0].font.name = font_name

        elif dec_type == "branding":
            text = dec.get("text", "")
            if not text:
                return
            color = self._resolve_color(dec.get("color", "text_light"))
            font_size = dec.get("font_size", 9)
            position = dec.get("position", "top_right")
            slide_w = styles.get("slide_width", 13.33)
            if position == "top_right":
                left = Inches(slide_w - 2.0)
                top = Inches(0.2)
            else:
                left = Inches(0.5)
                top = Inches(0.2)
            txBox = slide.shapes.add_textbox(left, top, Inches(2.0), Inches(0.3))
            tf = txBox.text_frame
            tf.paragraphs[0].text = text
            tf.paragraphs[0].font.size = Pt(font_size)
            tf.paragraphs[0].font.color.rgb = self._rgb(color)
            tf.paragraphs[0].font.name = font_name
            tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
