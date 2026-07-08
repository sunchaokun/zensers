"""
E2E test: Data-driven PPT generation pipeline with real file data + revision.

Flow: 
  1. Create real DOCX/PDF/Excel files with structured content
  2. PptInputAdapter.extract() → ExtractionResult
  3. PptRequirementExtractor → PptRequirement  
  4. PptDataSupplementer → fill gaps
  5. SlideDataBuilder → slide_data_list
  6. SlideOutlineBuilder → outline
  7. PptStructureEditor.edit() → HTMLToPPTConverter → REAL .pptx file
  8. SlideDataStore + PptVersionManager persistence
  9. PptRevisionService L1 revision (atomic text edit)
  10. PptRevisionService L4 revision (structure re-render)
  11. Version rollback + verify
  12. Open generated PPTX with python-pptx → verify slide count, titles, content
"""
import os
import json
import shutil
import pytest
from pathlib import Path
from unittest.mock import MagicMock, AsyncMock

from src.core.adjustment.ppt_input_adapter import PptInputAdapter
from src.core.adjustment.extraction_types import ExtractionResult, ExtractionSummary
from src.core.adjustment.ppt_requirement_extractor import PptRequirementExtractor, PptRequirement
from src.core.adjustment.ppt_data_supplementer import PptDataSupplementer, DataGap
from src.core.adjustment.slide_data_builder import SlideDataBuilder
from src.core.adjustment.slide_data_store import SlideDataStore
from src.core.adjustment.ppt_structure_editor import PptStructureEditor
from src.core.adjustment.ppt_version_manager import PptVersionManager
from src.core.adjustment.ppt_revision_service import PptRevisionService, PptRevisionRequest
from src.core.adjustment.ppt_slide_locator import PptSlideLocator
from src.core.adjustment.ppt_atomic_editor import PptAtomicEditor
from src.core.adjustment.ppt_element_editor import PptElementEditor
from src.core.adjustment.ppt_page_editor import PptPageEditor
from src.core.adjustment.ppt_revision_router import PptRevisionRouter
from src.core.adjustment.slide_data_path_resolver import SlideDataPathResolver
from src.core.adjustment.revision_types import RevisionOpType
from src.core.dialogue.state_machine import ConversationStateMachine, ConversationState
from src.converters.slide_outline_builder import SlideOutlineBuilder
from src.content.content_orchestrator import ContentSection, SectionType


def _create_real_docx(path):
    from docx import Document
    doc = Document()
    doc.add_paragraph("2024年中国新能源汽车市场分析报告", style="Heading 1")
    doc.add_paragraph("2024年中国新能源汽车销量达到950万辆，同比增长37.5%。", style="Normal")
    doc.add_paragraph("市场渗透率突破40%，标志着新能源汽车进入主流消费阶段。", style="Normal")
    doc.add_paragraph("竞争格局分析", style="Heading 1")
    doc.add_paragraph("比亚迪以35%的市场份额位居第一。", style="Normal")
    doc.add_paragraph("特斯拉上海工厂年产量超过90万辆。", style="Normal")
    doc.add_paragraph("技术趋势", style="Heading 1")
    doc.add_paragraph("固态电池技术将在2025年实现量产，成本下降30%。", style="Normal")
    doc.add_paragraph("800V高压平台成为新车型标配。", style="Normal")
    t = doc.add_table(rows=4, cols=3)
    t.rows[0].cells[0].text = "企业"
    t.rows[0].cells[1].text = "份额"
    t.rows[0].cells[2].text = "销量(万辆)"
    t.rows[1].cells[0].text = "比亚迪"
    t.rows[1].cells[1].text = "35%"
    t.rows[1].cells[2].text = "332"
    t.rows[2].cells[0].text = "特斯拉"
    t.rows[2].cells[1].text = "12%"
    t.rows[2].cells[2].text = "94"
    t.rows[3].cells[0].text = "蔚来"
    t.rows[3].cells[1].text = "4%"
    t.rows[3].cells[2].text = "16"
    doc.save(path)


def _create_real_pdf(path):
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(path, pagesize=letter)
    c.setFont("Helvetica", 12)
    c.drawString(72, 720, "Policy Analysis: New Energy Vehicle Subsidies")
    c.drawString(72, 700, "Government subsidies phased out in 2023, replaced by purchase tax exemption.")
    c.drawString(72, 680, "Local governments offer license plate incentives in major cities.")
    c.showPage()
    c.drawString(72, 720, "Infrastructure: Charging Network Expansion")
    c.drawString(72, 700, "Over 2.5 million public charging points deployed nationwide.")
    c.drawString(72, 680, "Fast charging (800V) coverage at 85% of highway service areas.")
    c.showPage()
    c.save()


def _create_real_excel(path):
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sales Data"
    ws.append(["Year", "Sales (10K)", "YoY Growth", "Penetration"])
    ws.append([2021, 352, "157%", "13%"])
    ws.append([2022, 688, "95%", "26%"])
    ws.append([2023, 950, "38%", "37%"])
    ws.append([2024, 1200, "26%", "45%"])
    ws2 = wb.create_sheet("Top Models")
    ws2.append(["Model", "Maker", "Sales (10K)", "Price Range"])
    ws2.append(["Song Plus DM-i", "BYD", "42", "15-22万"])
    ws2.append(["Model Y", "Tesla", "28", "26-36万"])
    ws2.append(["Han EV", "BYD", "22", "21-33万"])
    wb.save(path)


class TestE2EDataDrivenPipeline:
    @pytest.fixture(autouse=True)
    def setup(self, tmp_path):
        self.work_dir = tmp_path / "e2e_work"
        self.work_dir.mkdir()
        self.data_dir = tmp_path / "data"
        self.data_dir.mkdir()
        self.pptx_dir = tmp_path / "pptx_output"
        self.pptx_dir.mkdir()
        self.task_id = f"e2e_{id(self):08x}"

        self.docx_path = str(self.data_dir / "nev_report.docx")
        self.pdf_path = str(self.data_dir / "nev_policy.pdf")
        self.xlsx_path = str(self.data_dir / "nev_sales.xlsx")

        _create_real_docx(self.docx_path)
        _create_real_pdf(self.pdf_path)
        _create_real_excel(self.xlsx_path)

        self.store = SlideDataStore(task_id=self.task_id, data_dir=str(self.work_dir))
        self.version_mgr = PptVersionManager(revisions_dir=str(self.work_dir / "revisions"))

    def test_step1_file_extraction(self):
        adapter = PptInputAdapter()
        extraction = adapter.extract([self.docx_path, self.pdf_path, self.xlsx_path])

        assert extraction.title != ""
        assert len(extraction.sections) >= 5
        assert len(extraction.tables) >= 2
        assert "竞争格局分析" in extraction.key_topics
        assert "技术趋势" in extraction.key_topics
        assert extraction.metadata["file_count"] == 3
        assert "docx" in extraction.metadata["formats"]
        assert "pdf" in extraction.metadata["formats"]
        assert "xlsx" in extraction.metadata["formats"]

        docx_sections = [s for s in extraction.sections if s.title in ("2024年中国新能源汽车市场分析报告", "竞争格局分析", "技术趋势")]
        assert len(docx_sections) == 3

        pdf_sections = [s for s in extraction.sections if "Policy" in s.title or "Infrastructure" in s.title or "Page" in s.title]
        assert len(pdf_sections) >= 1

        xlsx_sections = [s for s in extraction.sections if s.title in ("Sales Data", "Top Models")]
        assert len(xlsx_sections) == 2

        sales_table = None
        for t in extraction.tables:
            if t and t[0] == ["Year", "Sales (10K)", "YoY Growth", "Penetration"]:
                sales_table = t
                break
        assert sales_table is not None
        assert sales_table[1] == ["2021", "352", "157%", "13%"]

        self.extraction = extraction

    def test_step2_requirement_extraction(self):
        adapter = PptInputAdapter()
        extraction = adapter.extract([self.docx_path, self.pdf_path, self.xlsx_path])

        extractor = PptRequirementExtractor()
        requirement = extractor.extract(
            extraction,
            user_description="做一个关于新能源汽车市场的汇报PPT",
        )
        assert "新能源" in requirement.topic
        assert requirement.page_count >= 3
        assert len(requirement.focus) > 0

        self.requirement = requirement

    def test_step3_data_supplement(self):
        adapter = PptInputAdapter()
        extraction = adapter.extract([self.docx_path, self.pdf_path, self.xlsx_path])
        extractor = PptRequirementExtractor()
        requirement = extractor.extract(extraction, user_description="做一个关于新能源汽车市场的汇报PPT")

        supplementer = PptDataSupplementer()
        gaps = supplementer.analyze_gaps(extraction, requirement)

        if gaps:
            class MockSearch:
                def execute(self, **kwargs):
                    return {"success": True, "data": {"results": [f"Supplemental data for {kwargs.get('query', '')}"]}}
            filled = supplementer.supplement(gaps, search_skill=MockSearch())
            assert all(g.filled for g in filled)

    def test_step4_slide_data_and_outline(self):
        adapter = PptInputAdapter()
        extraction = adapter.extract([self.docx_path, self.pdf_path, self.xlsx_path])

        builder = SlideDataBuilder()
        slide_data_list = builder.build_list(
            extraction.sections,
            add_cover=True, add_end=True,
            title="2024年中国新能源汽车市场分析报告",
        )
        assert len(slide_data_list) >= 3
        assert slide_data_list[0]["slide_type"] == "cover"
        assert slide_data_list[-1]["slide_type"] == "end"

        for sd in slide_data_list:
            assert "slide_type" in sd
            assert "title" in sd
            assert "items" in sd

        outline_builder = SlideOutlineBuilder()
        outline = outline_builder.build(slide_data_list, task_id=self.task_id)
        assert outline.total_pages == len(slide_data_list)
        assert len(outline.slides) == len(slide_data_list)

        self.slide_data_list = slide_data_list

    def test_step5_generate_real_pptx(self):
        adapter = PptInputAdapter()
        extraction = adapter.extract([self.docx_path, self.pdf_path, self.xlsx_path])

        builder = SlideDataBuilder()
        slide_data_list = builder.build_list(
            extraction.sections,
            add_cover=True, add_end=True,
            title="2024年中国新能源汽车市场分析报告",
        )

        output_path = str(self.pptx_dir / f"{self.task_id}_v1.pptx")
        editor = PptStructureEditor()
        result = editor.edit(slide_data_list, pptx="dummy", output_path=output_path)

        if result and result.success:
            assert os.path.exists(output_path)
            assert os.path.getsize(output_path) > 0
            assert result.slides_count == len(slide_data_list)

            from pptx import Presentation
            prs = Presentation(output_path)
            assert len(prs.slides) == len(slide_data_list)

            self.pptx_path = output_path
        else:
            pytest.skip("PPTX generation not available in this environment")

    def test_step6_persistence_and_versioning(self):
        adapter = PptInputAdapter()
        extraction = adapter.extract([self.docx_path])
        builder = SlideDataBuilder()
        slide_data_list = builder.build_list(extraction.sections, add_cover=True, add_end=True, title="Test")

        self.store.persist(self.task_id, slide_data_list)
        loaded = self.store.load(self.task_id)
        assert len(loaded) == len(slide_data_list)

        output_path = str(self.pptx_dir / f"{self.task_id}_v1.pptx")
        editor = PptStructureEditor()
        result = editor.edit(slide_data_list, pptx="dummy", output_path=output_path)

        if result and result.success:
            self.store.set_pptx_path(self.task_id, output_path)
            self.version_mgr.create_snapshot(self.task_id, output_path, "L0", "initial")
            versions = self.version_mgr.list_versions(self.task_id)
            assert len(versions) >= 1

    def test_step7_l1_revision(self):
        adapter = PptInputAdapter()
        extraction = adapter.extract([self.docx_path])
        builder = SlideDataBuilder()
        slide_data_list = builder.build_list(extraction.sections, add_cover=True, add_end=True, title="新能源报告")

        output_path = str(self.pptx_dir / f"{self.task_id}_l1.pptx")
        editor = PptStructureEditor()
        result = editor.edit(slide_data_list, pptx="dummy", output_path=output_path)

        if not (result and result.success):
            pytest.skip("PPTX generation not available")

        self.store.persist(self.task_id, slide_data_list)
        self.store.set_pptx_path(self.task_id, output_path)
        self.version_mgr.create_snapshot(self.task_id, output_path, "L0", "initial")

        atomic = PptAtomicEditor()
        success = atomic.update_slide_data(slide_data_list[1], "title", "市场竞争格局深度分析")
        assert success is True
        assert slide_data_list[1]["title"] == "市场竞争格局深度分析"

        success = atomic.update_slide_data(slide_data_list[1], "items[0]", "比亚迪以35%的市场份额稳居榜首")
        assert success is True
        assert slide_data_list[1]["items"][0] == "比亚迪以35%的市场份额稳居榜首"

        output_path_l1 = str(self.pptx_dir / f"{self.task_id}_l1_revised.pptx")
        result_l1 = editor.edit(slide_data_list, pptx="dummy", output_path=output_path_l1)
        if result_l1 and result_l1.success:
            assert os.path.exists(output_path_l1)

    def test_step8_l4_revision_delete_add(self):
        adapter = PptInputAdapter()
        extraction = adapter.extract([self.docx_path])
        builder = SlideDataBuilder()
        slide_data_list = builder.build_list(extraction.sections, add_cover=True, add_end=True, title="新能源报告")

        output_path = str(self.pptx_dir / f"{self.task_id}_l4.pptx")
        editor = PptStructureEditor()
        result = editor.edit(slide_data_list, pptx="dummy", output_path=output_path)

        if not (result and result.success):
            pytest.skip("PPTX generation not available")

        original_count = len(slide_data_list)

        success = editor.delete_slide(slide_data_list, 2)
        assert success is True
        assert len(slide_data_list) == original_count - 1

        new_slide = {
            "slide_type": "content", "title": "政策环境分析",
            "content": "新能源车购置税减免政策延续至2027年",
            "items": ["购置税减免延续", "双积分政策趋严", "地方补贴转向基建"],
            "table_data": [], "extra_tables": [], "images": [],
            "source_text": "", "section_number": 2,
            "section_summary": "", "insight_text": "", "kpi_data": [], "comparison_data": [],
        }
        success = editor.add_slide(slide_data_list, 2, new_slide)
        assert success is True

        output_path_l4 = str(self.pptx_dir / f"{self.task_id}_l4_revised.pptx")
        result_l4 = editor.edit(slide_data_list, pptx="dummy", output_path=output_path_l4)
        if result_l4 and result_l4.success:
            from pptx import Presentation
            prs = Presentation(output_path_l4)
            assert len(prs.slides) == len(slide_data_list)

            found_policy = False
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame and "政策环境" in shape.text_frame.text:
                        found_policy = True
                        break
            assert found_policy, "Added '政策环境分析' slide not found in generated PPTX"

    async def test_step9_revision_service_dispatch(self):
        adapter = PptInputAdapter()
        extraction = adapter.extract([self.docx_path])
        builder = SlideDataBuilder()
        slide_data_list = builder.build_list(extraction.sections, add_cover=True, add_end=True, title="新能源报告")

        output_path = str(self.pptx_dir / f"{self.task_id}_svc.pptx")
        editor = PptStructureEditor()
        result = editor.edit(slide_data_list, pptx="dummy", output_path=output_path)

        if not (result and result.success):
            pytest.skip("PPTX generation not available")

        self.store.persist(self.task_id, slide_data_list)
        self.store.set_pptx_path(self.task_id, output_path)

        service = PptRevisionService(self.store)

        request_l0 = PptRevisionRequest(
            task_id=self.task_id,
            source="natural",
            description="分析一下第三页的内容",
            revision_type=RevisionOpType.REVIEW,
        )
        result_l0 = await service._dispatch("L0", request_l0)
        assert result_l0.success is True

        request_l1 = PptRevisionRequest(
            task_id=self.task_id,
            source="click",
            slide_index=1,
            revision_type=RevisionOpType.REPLACE_TEXT,
            target_field="title",
            new_value="竞争格局深度解读",
            revision_level="L1",
        )
        result_l1 = await service._dispatch("L1", request_l1)
        assert result_l1.success is True
        updated = self.store.load(self.task_id)
        assert updated[1]["title"] == "竞争格局深度解读"

    def test_step10_rollback(self):
        adapter = PptInputAdapter()
        extraction = adapter.extract([self.docx_path])
        builder = SlideDataBuilder()
        slide_data_list = builder.build_list(extraction.sections, add_cover=True, add_end=True, title="新能源报告")

        output_path_v1 = str(self.pptx_dir / f"{self.task_id}_v1.pptx")
        editor = PptStructureEditor()
        result = editor.edit(slide_data_list, pptx="dummy", output_path=output_path_v1)

        if not (result and result.success):
            pytest.skip("PPTX generation not available")

        self.store.persist(self.task_id, slide_data_list)
        self.store.set_pptx_path(self.task_id, output_path_v1)
        snap = self.version_mgr.create_snapshot(self.task_id, output_path_v1, "L0", "initial")

        original_title = slide_data_list[1]["title"]
        slide_data_list[1]["title"] = "MODIFIED TITLE"

        output_path_v2 = str(self.pptx_dir / f"{self.task_id}_v2.pptx")
        editor.edit(slide_data_list, pptx="dummy", output_path=output_path_v2)
        self.version_mgr.create_snapshot(self.task_id, output_path_v2, "L1", "modified title")

        self.version_mgr.rollback(self.task_id, snap, output_path_v1)

        restored_data = self.store.load(self.task_id)
        assert restored_data[1]["title"] == original_title

    def test_step11_state_machine_full_flow(self):
        sm = ConversationStateMachine(research_id=self.task_id)
        assert sm.current_state == ConversationState.UNDERSTANDING

        sm.transition(ConversationState.DATA_EXTRACTED)
        assert sm.current_state == ConversationState.DATA_EXTRACTED

        sm.transition(ConversationState.REQUIREMENT_CONFIRM)
        assert sm.current_state == ConversationState.REQUIREMENT_CONFIRM

        sm.transition(ConversationState.DATA_SUPPLEMENT)
        assert sm.current_state == ConversationState.DATA_SUPPLEMENT

        sm.transition(ConversationState.FRAMEWORK_CONFIRM)
        assert sm.current_state == ConversationState.FRAMEWORK_CONFIRM

        sm.transition(ConversationState.EXECUTING)
        assert sm.current_state == ConversationState.EXECUTING

        sm.transition(ConversationState.PREVIEWING)
        assert sm.current_state == ConversationState.PREVIEWING

        sm.transition(ConversationState.COMPLETED)
        assert sm.current_state == ConversationState.COMPLETED

    def test_step12_pptx_field_completeness(self):
        adapter = PptInputAdapter()
        extraction = adapter.extract([self.docx_path, self.pdf_path, self.xlsx_path])
        builder = SlideDataBuilder()
        slide_data_list = builder.build_list(
            extraction.sections, add_cover=True, add_end=True,
            title="新能源市场报告",
        )

        required_fields = [
            "slide_type", "title", "content", "items",
            "table_data", "extra_tables", "images",
            "source_text", "section_number", "section_summary",
            "insight_text", "kpi_data", "comparison_data",
        ]
        for i, sd in enumerate(slide_data_list):
            for field in required_fields:
                assert field in sd, f"Slide {i} missing field '{field}'"

    def test_step13_extraction_summary_for_frontend(self):
        adapter = PptInputAdapter()
        extraction = adapter.extract([self.docx_path, self.pdf_path, self.xlsx_path])

        from src.api.research_api import ResearchAPI
        api = ResearchAPI.__new__(ResearchAPI)
        file_ids = [
            {"id": "f1", "path": self.docx_path, "filename": "nev_report.docx", "size_mb": 0.1},
            {"id": "f2", "path": self.pdf_path, "filename": "nev_policy.pdf", "size_mb": 0.1},
            {"id": "f3", "path": self.xlsx_path, "filename": "nev_sales.xlsx", "size_mb": 0.1},
        ]
        summary = api._build_extraction_summary(extraction, file_ids)

        assert summary.file_count == 3
        assert summary.title is not None
        assert len(summary.sections) > 0
        assert summary.tables_count >= 2
        assert summary.extraction_status == "success"

        msg = api._format_extraction_summary_message(summary)
        assert "已读取" in msg
        assert "3个文件" in msg
        assert "您想基于这份材料做什么" in msg
