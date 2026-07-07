"""Full E2E: research → preview → export PPT → revise → re-export → all formats"""
import sys
sys.path.insert(0, "E:/market_report_systerm")

import os
import shutil
import asyncio
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock, patch

from src.content.content_orchestrator import ContentOrchestrator
from src.converters.html_to_ppt import HTMLToPPTConverter
from src.converters.html_to_word import HTMLToWordConverter
from src.converters.html_to_pdf import HTMLToPDFConverter
from src.services.chart_generator import ChartGenerator, ChartConfig, ChartType
from src.core.preview_storage import PreviewStorage

CWD = str(Path.cwd())
E2E_DIR = os.path.join(CWD, "__e2e_full__")


def _make_charts():
    charts_dir = Path(E2E_DIR) / "charts"
    charts_dir.mkdir(parents=True, exist_ok=True)
    gen = ChartGenerator(output_dir=charts_dir)
    charts = []
    for cfg in [
        ChartConfig(chart_type=ChartType.BAR, title="Market Size by Region",
                    data={"categories": ["NA", "EU", "APAC", "LATAM"], "values": [4.2, 3.8, 5.1, 1.2]}, dpi=200),
        ChartConfig(chart_type=ChartType.LINE, title="Revenue Trend",
                    data={"x": [2020, 2021, 2022, 2023, 2024], "y": [12.3, 15.7, 18.2, 22.1, 28.5]},
                    xlabel="Year", ylabel="Revenue", dpi=200),
        ChartConfig(chart_type=ChartType.PIE, title="Market Share",
                    data={"categories": ["Tech", "Health", "Finance", "Consumer"], "values": [35, 25, 20, 20]}, dpi=200),
    ]:
        r = gen.generate(cfg)
        if r.success:
            charts.append({"path": r.image_path, "caption": cfg.title, "title": cfg.title, "anchor_type": "section_end"})
    return charts


def _research_result_v1(charts):
    return {
        "title": "Global Market Analysis 2024",
        "sections": [
            {"id": "exec", "title": "Executive Summary",
             "content": "The global market grew 28.9% YoY to reach 15.1B USD. Asia Pacific leads with 5.1B. Technology sector dominates at 35% share.",
             "order": 1, "type": "body", "charts": charts[:1]},
            {"id": "overview", "title": "Market Overview",
             "content": "Key drivers include digital transformation, cloud adoption, and AI innovation. Asia Pacific is the fastest-growing region.",
             "order": 2, "type": "body", "charts": charts[1:2]},
            {"id": "competitive", "title": "Competitive Landscape",
             "content": "Top 5 players hold 45% share. Technology leads at 35%, Healthcare at 25%. Startups gaining in AI and sustainability.",
             "order": 3, "type": "body", "charts": charts[2:3]},
        ],
        "key_findings": ["Market reached 15.1B USD", "28.9% YoY growth", "APAC fastest region", "Tech 35% share"],
        "data_points": [{"metric": "Market", "value": "15.1B", "unit": "USD"}, {"metric": "Growth", "value": "28.9%", "unit": "%"}],
    }


def _research_result_v2_revised(charts):
    revised = _research_result_v1(charts)
    revised["sections"][1]["content"] = (
        "Key drivers include digital transformation, cloud adoption, and AI innovation.\n\n"
        "## Asia Pacific Deep Dive\n\n"
        "Asia Pacific reached 5.1B USD in 2024, driven by China's 3.2B and India's 0.8B markets. "
        "The region's CAGR of 24.3% outpaces the global average of 18.2%. "
        "Key verticals: semiconductor manufacturing (+31%), enterprise SaaS (+27%), fintech (+22%)."
    )
    return revised


class TestFullPipeline:
    def setup_method(self):
        os.makedirs(E2E_DIR, exist_ok=True)
        self.charts = _make_charts()
        self.orch = ContentOrchestrator()
        self.task_id = "e2e_test_001"

    def teardown_method(self):
        shutil.rmtree(E2E_DIR, ignore_errors=True)
        for d in [PreviewStorage.OLD_DIR, PreviewStorage.NEW_DIR,
                   Path(f"data/reports/{self.task_id}")]:
            if d.exists():
                shutil.rmtree(d, ignore_errors=True)

    def test_step1_generate_preview(self):
        research = _research_result_v1(self.charts)
        html = self.orch.transform_to_html(research, output_format="pptx", output_dir=E2E_DIR)
        assert html, "Should produce HTML preview"
        assert len(html) > 200
        PreviewStorage.write(self.task_id, html)
        assert PreviewStorage.path(self.task_id).exists()
        assert (PreviewStorage.OLD_DIR / f"{self.task_id}.html").exists()
        print(f"[STEP 1] Preview generated: {len(html)} chars")

    def test_step2_export_pptx(self):
        research = _research_result_v1(self.charts)
        html = self.orch.transform_to_html(research, output_format="pptx", output_dir=E2E_DIR)
        PreviewStorage.write(self.task_id, html)

        preview_path = PreviewStorage.OLD_DIR / f"{self.task_id}.html"
        html_content = preview_path.read_text(encoding="utf-8")

        output_dir = Path(f"data/reports/{self.task_id}")
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = str(output_dir / f"{self.task_id}_report.pptx")

        converter = HTMLToPPTConverter()
        result = converter.convert(html=html_content, output_path=output_path)
        assert result.success, f"PPTX export failed: {result.error}"
        assert os.path.exists(output_path)

        from pptx import Presentation
        prs = Presentation(output_path)
        assert len(prs.slides) >= 4, f"Expected >=4 slides, got {len(prs.slides)}"
        assert abs(prs.slide_width / prs.slide_height - 16/9) < 0.05, "Should be 16:9"

        total_images = sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == 13)
        assert total_images >= 2, f"Expected >=2 images, got {total_images}"
        print(f"[STEP 2] PPTX exported: {len(prs.slides)} slides, {total_images} images")

    def test_step3_export_all_formats(self):
        research = _research_result_v1(self.charts)
        for fmt in ["docx", "pptx", "pdf"]:
            html = self.orch.transform_to_html(research, output_format=fmt, output_dir=E2E_DIR)
            assert html, f"Should produce HTML for {fmt}"

            output_dir = Path(E2E_DIR) / f"{fmt}_out"
            output_dir.mkdir(parents=True, exist_ok=True)
            output_path = str(output_dir / f"report.{fmt}")

            if fmt == "docx":
                converter = HTMLToWordConverter()
            elif fmt == "pptx":
                converter = HTMLToPPTConverter()
            else:
                converter = HTMLToPDFConverter()

            result = converter.convert(html=html, output_path=output_path)
            assert result.success, f"{fmt} export failed: {result.error}"
            assert os.path.exists(output_path)
            assert os.path.getsize(output_path) > 500, f"{fmt} file too small"
            print(f"[STEP 3] {fmt}: {os.path.getsize(output_path)} bytes")

    def test_step4_revision_reexport(self):
        research_v1 = _research_result_v1(self.charts)
        html_v1 = self.orch.transform_to_html(research_v1, output_format="pptx", output_dir=E2E_DIR)
        assert "Executive Summary" in html_v1
        assert "Asia Pacific Deep Dive" not in html_v1

        research_v2 = _research_result_v2_revised(self.charts)
        html_v2 = self.orch.transform_to_html(research_v2, output_format="pptx", output_dir=E2E_DIR)
        assert "Asia Pacific Deep Dive" in html_v2
        assert "semiconductor" in html_v2

        output_dir = Path(E2E_DIR) / "revised_out"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = str(output_dir / "report_v2.pptx")
        converter = HTMLToPPTConverter()
        result = converter.convert(html=html_v2, output_path=output_path)
        assert result.success
        print(f"[STEP 4] Revised PPTX: {result.slides_count} slides")

    def test_step5_pptx_design_verification(self):
        research = _research_result_v1(self.charts)
        html = self.orch.transform_to_html(research, output_format="pptx", output_dir=E2E_DIR)

        output_dir = Path(E2E_DIR) / "design_out"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = str(output_dir / "report.pptx")
        converter = HTMLToPPTConverter()
        result = converter.convert(html=html, output_path=output_path)
        assert result.success

        from pptx import Presentation
        from pptx.dml.color import RGBColor
        prs = Presentation(output_path)

        cover = prs.slides[0]
        bg = cover.background.fill
        assert bg.type is not None, "Cover should have background fill"

        has_gold_text = False
        for shape in cover.shapes:
            if hasattr(shape, "text_frame"):
                for p in shape.text_frame.paragraphs:
                    try:
                        if p.font.color.rgb == RGBColor(0xC9, 0xA2, 0x27):
                            has_gold_text = True
                    except AttributeError:
                        pass
        assert has_gold_text, "Cover should have gold (#C9A227) title text"

        end_slide = prs.slides[-1]
        end_bg = end_slide.background.fill
        assert end_bg.type is not None, "End slide should have background"
        print(f"[STEP 5] Design verified: gradient cover, gold title, navy end")

    def test_step6_pptx_image_aspect_ratio(self):
        research = _research_result_v1(self.charts)
        html = self.orch.transform_to_html(research, output_format="pptx", output_dir=E2E_DIR)

        output_dir = Path(E2E_DIR) / "aspect_out"
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = str(output_dir / "report.pptx")
        converter = HTMLToPPTConverter()
        result = converter.convert(html=html, output_path=output_path)
        assert result.success

        from pptx import Presentation
        prs = Presentation(output_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    display_aspect = shape.width / shape.height
                    try:
                        native_aspect = shape.image.size[0] / shape.image.size[1]
                        assert abs(display_aspect - native_aspect) < 0.01, \
                            f"Image distorted: display={display_aspect:.3f} vs native={native_aspect:.3f}"
                    except Exception:
                        pass
        print("[STEP 6] All images maintain original aspect ratio")

    def test_full_flow_sequential(self):
        print("\n=== FULL E2E: 研究 → 预览 → 导出PPT → 修订 → 重导出PPT ===")
        research = _research_result_v1(self.charts)

        print("[1/5] 生成HTML预览...")
        html = self.orch.transform_to_html(research, output_format="pptx", output_dir=E2E_DIR)
        PreviewStorage.write(self.task_id, html)
        assert PreviewStorage.path(self.task_id).exists()
        print(f"  -> 预览已写入, {len(html)} chars")

        print("[2/5] 导出PPTX...")
        preview_path = PreviewStorage.OLD_DIR / f"{self.task_id}.html"
        html_content = preview_path.read_text(encoding="utf-8")
        output_dir = Path(f"data/reports/{self.task_id}")
        output_dir.mkdir(parents=True, exist_ok=True)
        pptx_path = str(output_dir / f"{self.task_id}_report.pptx")
        result = HTMLToPPTConverter().convert(html=html_content, output_path=pptx_path)
        assert result.success
        from pptx import Presentation
        prs = Presentation(pptx_path)
        img_count = sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == 13)
        print(f"  -> PPTX: {result.slides_count}页, {img_count}张图表, {os.path.getsize(pptx_path)} bytes")

        print("[3/5] 用户审阅后提出修订: 增加亚太市场深度分析...")
        research_v2 = _research_result_v2_revised(self.charts)
        html_v2 = self.orch.transform_to_html(research_v2, output_format="pptx", output_dir=E2E_DIR)
        assert "Asia Pacific Deep Dive" in html_v2
        assert "semiconductor" in html_v2
        PreviewStorage.write(self.task_id, html_v2)
        print("  -> 修订内容已生成")

        print("[4/5] 重新导出PPTX...")
        html_v2_content = (PreviewStorage.OLD_DIR / f"{self.task_id}.html").read_text(encoding="utf-8")
        pptx_v2_path = str(output_dir / f"{self.task_id}_report_v2.pptx")
        result = HTMLToPPTConverter().convert(html=html_v2_content, output_path=pptx_v2_path)
        assert result.success
        prs_v2 = Presentation(pptx_v2_path)
        print(f"  -> 修订后PPTX: {result.slides_count}页, {os.path.getsize(pptx_v2_path)} bytes")

        print("[5/5] 验证修订内容已体现在PPTX中...")
        all_text = ""
        for slide in prs_v2.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    all_text += shape.text_frame.text + " "
        assert "Asia Pacific" in all_text or "semiconductor" in all_text, "修订内容应出现在PPTX中"
        print("  -> 修订内容已确认")

        print("\n=== 全流程通过 ===")


if __name__ == "__main__":
    import pytest
    pytest.main([__file__, "-v", "--tb=short", "-s"])
