"""End-to-end tests: full pipeline from ContentOrchestrator to DOCX/PPTX/PDF"""
import sys
sys.path.insert(0, "E:/market_report_systerm")

import os
import shutil
from pathlib import Path

from src.content.content_orchestrator import ContentOrchestrator, ContentSection, SectionType
from src.converters.html_to_word import HTMLToWordConverter
from src.converters.html_to_ppt import HTMLToPPTConverter
from src.converters.html_to_pdf import HTMLToPDFConverter
from src.services.chart_generator import ChartGenerator, ChartConfig, ChartType
from src.content.format_strategy import get_format_strategy, FormatStrategy

CWD = str(Path.cwd())
E2E_TMP = os.path.join(CWD, "__e2e_test_tmp__")


def _generate_real_charts(output_dir: str):
    """Generate real chart PNGs using ChartGenerator."""
    gen = ChartGenerator(output_dir=Path(output_dir))
    
    charts = []
    
    bar_config = ChartConfig(
        chart_type=ChartType.BAR,
        title="Global Market Size by Region (2024)",
        data={
            "categories": ["North America", "Europe", "Asia Pacific", "Latin America", "Middle East"],
            "values": [4.2, 3.8, 5.1, 1.2, 0.8],
        },
        source="Market Research Institute",
        dpi=200,
    )
    result = gen.generate(bar_config)
    if result.success:
        charts.append({"path": result.image_path, "caption": "Global Market Size", "title": "Market Size", "anchor_type": "section_end"})
    
    line_config = ChartConfig(
        chart_type=ChartType.LINE,
        title="Revenue Growth Trend (2020-2024)",
        data={
            "x": [2020, 2021, 2022, 2023, 2024],
            "y": [12.3, 15.7, 18.2, 22.1, 28.5],
        },
        xlabel="Year",
        ylabel="Revenue (Billion USD)",
        dpi=200,
    )
    result = gen.generate(line_config)
    if result.success:
        charts.append({"path": result.image_path, "caption": "Revenue Growth", "title": "Revenue Trend", "anchor_type": "section_end"})
    
    pie_config = ChartConfig(
        chart_type=ChartType.PIE,
        title="Market Share by Segment",
        data={
            "categories": ["Technology", "Healthcare", "Finance", "Consumer", "Industrial"],
            "values": [35, 25, 20, 12, 8],
        },
        dpi=200,
    )
    result = gen.generate(pie_config)
    if result.success:
        charts.append({"path": result.image_path, "caption": "Market Share", "title": "Segment Share", "anchor_type": "section_end"})
    
    return charts


def _build_research_result(charts: list) -> dict:
    """Build a realistic research_result dict."""
    return {
        "title": "Global Market Analysis Report 2024",
        "sections": [
            {
                "id": "exec_summary",
                "title": "Executive Summary",
                "content": "The global market has experienced significant growth in 2024, driven by technological innovation and expanding consumer bases in emerging economies. Key findings include a 28.9% year-over-year revenue increase, with Asia Pacific leading growth at 5.1 billion USD. The technology sector continues to dominate with 35% market share.\n\nMarket headwinds include regulatory uncertainty in the EU and supply chain disruptions affecting the industrial segment. However, strong fundamentals and accelerating digital transform_to_htmlation suggest sustained growth through 2025.",
                "order": 1,
                "type": "executive_summary",
                "charts": charts[:1],
            },
            {
                "id": "market_overview",
                "title": "Market Overview and Size",
                "content": "The global market reached an estimated value of 15.1 billion USD in 2024, representing a compound annual growth rate (CAGR) of 18.2% since 2020.\n\n## Regional Performance\n\nAsia Pacific emerged as the fastest-growing region, driven by rapid industrialization and increasing technology adoption. North America maintained its position as the largest single market.\n\n## Key Growth Drivers\n\n- Digital transform_to_htmlation accelerating across industries\n- Cloud computing adoption reaching critical mass\n- AI and machine learning creating new market opportunities\n- Government initiatives supporting tech infrastructure",
                "order": 2,
                "type": "body",
                "charts": charts[1:2],
            },
            {
                "id": "competitive",
                "title": "Competitive Landscape",
                "content": "The competitive landscape is characterized by increasing consolidation and strategic partnerships. The top 5 players account for approximately 45% of total market share.\n\n## Technology Sector Dominance\n\nThe technology sector commands 35% market share, followed by Healthcare at 25%. This reflects the growing convergence of technology and healthcare, particularly in telemedicine and health-tech applications.\n\n## Emerging Competitors\n\nSeveral startups have gained significant market traction, particularly in the AI and sustainability segments. Traditional players are responding through acquisitions and internal innovation programs.",
                "order": 3,
                "type": "body",
                "charts": charts[2:3],
            },
        ],
        "key_findings": [
            "Global market reached 15.1B USD in 2024 (28.9% YoY growth)",
            "Asia Pacific is the fastest-growing region at 5.1B USD",
            "Technology sector leads with 35% market share",
            "CAGR of 18.2% projected to continue through 2025",
        ],
        "data_points": [
            {"metric": "Total Market Value", "value": "15.1B", "unit": "USD"},
            {"metric": "YoY Growth", "value": "28.9%", "unit": "%"},
            {"metric": "CAGR (2020-2024)", "value": "18.2%", "unit": "%"},
            {"metric": "Top Region", "value": "Asia Pacific", "unit": "-"},
        ],
    }


class TestE2EDocxPipeline:
    def setup_method(self):
        os.makedirs(E2E_TMP, exist_ok=True)
        self.charts = _generate_real_charts(os.path.join(E2E_TMP, "charts"))
        self.research = _build_research_result(self.charts)
        self.output_dir = os.path.join(E2E_TMP, "docx_output")
        os.makedirs(self.output_dir, exist_ok=True)
    
    def teardown_method(self):
        shutil.rmtree(E2E_TMP, ignore_errors=True)
    
    def test_full_docx_pipeline(self):
        orch = ContentOrchestrator()
        html = orch.transform_to_html(
            self.research, output_format="docx", output_dir=self.output_dir
        )
        assert html, "Orchestrator should produce HTML"
        assert len(html) > 100, "HTML should have substantial content"
        
        converter = HTMLToWordConverter()
        output_path = os.path.join(self.output_dir, "report.docx")
        result = converter.convert(html=html, output_path=output_path)
        assert result.success, f"DOCX conversion failed: {result.error}"
        assert os.path.exists(output_path), "Output file should exist"
        assert os.path.getsize(output_path) > 1000, "DOCX should have content"
        
        from docx import Document
        doc = Document(output_path)
        assert len(doc.paragraphs) > 0, "Should have paragraphs"


class TestE2EPptxPipeline:
    def setup_method(self):
        os.makedirs(E2E_TMP, exist_ok=True)
        self.charts = _generate_real_charts(os.path.join(E2E_TMP, "charts"))
        self.research = _build_research_result(self.charts)
        self.output_dir = os.path.join(E2E_TMP, "pptx_output")
        os.makedirs(self.output_dir, exist_ok=True)
    
    def teardown_method(self):
        shutil.rmtree(E2E_TMP, ignore_errors=True)
    
    def test_full_pptx_pipeline(self):
        orch = ContentOrchestrator()
        html = orch.transform_to_html(
            self.research, output_format="pptx", output_dir=self.output_dir
        )
        assert html, "Orchestrator should produce HTML for PPTX"
        assert 'class="slide' in html or 'data-type=' in html, "Should have slide sections"
        
        converter = HTMLToPPTConverter()
        output_path = os.path.join(self.output_dir, "report.pptx")
        result = converter.convert(html=html, output_path=output_path)
        assert result.success, f"PPTX conversion failed: {result.error}"
        assert result.slides_count and result.slides_count > 0, "Should have slides"
        
        from pptx import Presentation
        prs = Presentation(output_path)
        assert len(prs.slides) > 0, "Should have slides"
        
        w = prs.slide_width
        h = prs.slide_height
        from pptx.util import Inches
        ratio = w / h
        expected_ratio = 13.333 / 7.5
        assert abs(ratio - expected_ratio) < 0.05, f"Slide ratio should be ~16:9 ({expected_ratio}), got {ratio}"
    
    def test_pptx_contains_images(self):
        orch = ContentOrchestrator()
        html = orch.transform_to_html(
            self.research, output_format="pptx", output_dir=self.output_dir
        )
        
        converter = HTMLToPPTConverter()
        output_path = os.path.join(self.output_dir, "report_with_charts.pptx")
        result = converter.convert(html=html, output_path=output_path)
        assert result.success, f"PPTX conversion failed: {result.error}"
        
        from pptx import Presentation
        prs = Presentation(output_path)
        
        total_pictures = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    total_pictures += 1
        
        assert total_pictures > 0, f"Expected at least 1 image in PPTX, got {total_pictures}"


class TestE2EPdfPipeline:
    def setup_method(self):
        os.makedirs(E2E_TMP, exist_ok=True)
        self.charts = _generate_real_charts(os.path.join(E2E_TMP, "charts"))
        self.research = _build_research_result(self.charts)
        self.output_dir = os.path.join(E2E_TMP, "pdf_output")
        os.makedirs(self.output_dir, exist_ok=True)
    
    def teardown_method(self):
        shutil.rmtree(E2E_TMP, ignore_errors=True)
    
    def test_full_pdf_pipeline(self):
        orch = ContentOrchestrator()
        html = orch.transform_to_html(
            self.research, output_format="pdf", output_dir=self.output_dir
        )
        assert html, "Orchestrator should produce HTML for PDF"
        
        converter = HTMLToPDFConverter()
        output_path = os.path.join(self.output_dir, "report.pdf")
        result = converter.convert(html=html, output_path=output_path)
        assert result.success, f"PDF conversion failed: {result.error}"
        assert os.path.exists(output_path), "Output file should exist"
        assert os.path.getsize(output_path) > 5000, "PDF should have substantial content"
    
    def test_pdf_with_images_is_larger(self):
        orch = ContentOrchestrator()
        html_with_charts = orch.transform_to_html(
            self.research, output_format="pdf", output_dir=self.output_dir
        )
        
        html_no_charts = orch.transform_to_html(
            {k: v for k, v in self.research.items() if k != "sections"} | {
                "sections": [{k: v for k, v in s.items() if k != "charts"} for s in self.research["sections"]]
            },
            output_format="pdf", output_dir=self.output_dir
        )
        
        converter = HTMLToPDFConverter()
        
        path_with = os.path.join(self.output_dir, "with_charts.pdf")
        result_with = converter.convert(html=html_with_charts, output_path=path_with)
        assert result_with.success
        
        path_without = os.path.join(self.output_dir, "without_charts.pdf")
        result_without = converter.convert(html=html_no_charts, output_path=path_without)
        assert result_without.success
        
        size_with = os.path.getsize(path_with)
        size_without = os.path.getsize(path_without)
        assert size_with > size_without, f"PDF with charts ({size_with}) should be larger than without ({size_without})"


class TestE2EFormatStrategyIntegration:
    """Verify FormatStrategy configs match real converter behavior."""
    
    def test_pptx_chart_dpi_matches_strategy(self):
        strategy = get_format_strategy("pptx")
        style = strategy.get_chart_style()
        assert style.dpi == 200, "PPTX strategy should specify 200 DPI"
        assert style.transparent_bg is True, "PPTX charts should have transparent background"
    
    def test_docx_chart_dpi_matches_strategy(self):
        strategy = get_format_strategy("docx")
        style = strategy.get_chart_style()
        assert style.dpi == 150, "DOCX strategy should specify 150 DPI"
        assert style.transparent_bg is False, "DOCX charts should not be transparent"
    
    def test_pdf_chart_dpi_matches_strategy(self):
        strategy = get_format_strategy("pdf")
        style = strategy.get_chart_style()
        assert style.dpi == 200, "PDF strategy should specify 200 DPI"
        assert style.transparent_bg is False, "PDF charts should not be transparent"


class TestE2EMultiFormatConsistency:
    """Verify same research data produces valid output in all formats."""
    
    def setup_method(self):
        os.makedirs(E2E_TMP, exist_ok=True)
        self.charts = _generate_real_charts(os.path.join(E2E_TMP, "charts"))
        self.research = _build_research_result(self.charts)
    
    def teardown_method(self):
        shutil.rmtree(E2E_TMP, ignore_errors=True)
    
    def test_all_formats_produce_valid_output(self):
        orch = ContentOrchestrator()
        
        results = {}
        for fmt in ["docx", "pptx", "pdf"]:
            output_dir = os.path.join(E2E_TMP, f"{fmt}_output")
            os.makedirs(output_dir, exist_ok=True)
            
            html = orch.transform_to_html(
                self.research, output_format=fmt, output_dir=output_dir
            )
            assert html, f"Orchestrator should produce HTML for {fmt}"
            assert len(html) > 50, f"{fmt} HTML should have content"
            results[fmt] = html
        
        for fmt, html in results.items():
            output_dir = os.path.join(E2E_TMP, f"{fmt}_output")
            if fmt == "docx":
                converter = HTMLToWordConverter()
                output_path = os.path.join(output_dir, f"report.{fmt}")
            elif fmt == "pptx":
                converter = HTMLToPPTConverter()
                output_path = os.path.join(output_dir, f"report.{fmt}")
            else:
                converter = HTMLToPDFConverter()
                output_path = os.path.join(output_dir, f"report.{fmt}")
            
            result = converter.convert(html=html, output_path=output_path)
            assert result.success, f"{fmt} conversion failed: {result.error}"
            assert os.path.exists(output_path), f"{fmt} output file should exist"
            assert os.path.getsize(output_path) > 500, f"{fmt} file should have content"


if __name__ == "__main__":
    import pytest
    pytest.main([__file__, "-v", "--tb=short"])
