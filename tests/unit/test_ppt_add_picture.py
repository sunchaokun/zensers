"""Tests for HTMLToPPTConverter add_picture support (P0-1)"""
import sys
sys.path.insert(0, "E:/market_report_systerm")

import os
import tempfile
import shutil
from pathlib import Path

from src.converters.html_to_ppt import HTMLToPPTConverter, ConversionResult

CWD = str(Path.cwd())
TEST_TMP = os.path.join(CWD, "__ppt_test_tmp__")


class TestPPTImageSupport:
    def setup_method(self):
        os.makedirs(TEST_TMP, exist_ok=True)
        self.converter = HTMLToPPTConverter()
        self.output_path = os.path.join(TEST_TMP, "test.pptx")

        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        fig, ax = plt.subplots()
        ax.bar(["A", "B"], [1, 2])
        ax.set_title("Test Chart")
        self.chart_path = os.path.join(TEST_TMP, "chart.png")
        fig.savefig(self.chart_path, dpi=100)
        plt.close(fig)

    def teardown_method(self):
        if os.path.exists(TEST_TMP):
            shutil.rmtree(TEST_TMP, ignore_errors=True)

    def test_img_tag_in_html_produces_picture(self):
        html = (
            "<section class='slide' data-type='content'>"
            "<h1>Chart Slide</h1>"
            f"<img src='{self.chart_path}' alt='Test Chart'>"
            "</section>"
        )
        result = self.converter.convert(html, self.output_path)
        assert result.success, f"Conversion failed: {result.error}"

        from pptx import Presentation
        prs = Presentation(self.output_path)
        assert len(prs.slides) >= 1
        slide = prs.slides[0]
        has_picture = any(
            shape.shape_type == 13
            for shape in slide.shapes
        )
        assert has_picture, "Expected a picture shape in the slide but found none"

    def test_img_tag_with_text_and_image(self):
        html = (
            "<section class='slide' data-type='content'>"
            "<h1>Mixed Slide</h1>"
            "<p>Some text content</p>"
            f"<img src='{self.chart_path}' alt='Chart'>"
            "</section>"
        )
        result = self.converter.convert(html, self.output_path)
        assert result.success, f"Conversion failed: {result.error}"

        from pptx import Presentation
        prs = Presentation(self.output_path)
        slide = prs.slides[0]
        shape_types = [shape.shape_type for shape in slide.shapes]
        has_text = any(t in shape_types for t in [17, 14, 24])
        has_picture = 13 in shape_types
        assert has_text, "Expected text shape"
        assert has_picture, "Expected picture shape"

    def test_multiple_images_on_one_slide(self):
        chart2_path = os.path.join(TEST_TMP, "chart2.png")
        import matplotlib.pyplot as plt
        fig, ax = plt.subplots()
        ax.plot([1, 2, 3], [3, 1, 2])
        fig.savefig(chart2_path, dpi=100)
        plt.close(fig)

        html = (
            "<section class='slide' data-type='content'>"
            "<h1>Two Charts</h1>"
            f"<img src='{self.chart_path}' alt='Chart1'>"
            f"<img src='{chart2_path}' alt='Chart2'>"
            "</section>"
        )
        result = self.converter.convert(html, self.output_path)
        assert result.success, f"Conversion failed: {result.error}"

        from pptx import Presentation
        prs = Presentation(self.output_path)
        slide = prs.slides[0]
        picture_count = sum(1 for s in slide.shapes if s.shape_type == 13)
        assert picture_count == 2, f"Expected 2 pictures, got {picture_count}"

    def test_missing_image_src_does_not_crash(self):
        html = (
            "<section class='slide' data-type='content'>"
            "<h1>No Image</h1>"
            "<img src='/nonexistent/path.png' alt='Missing'>"
            "</section>"
        )
        result = self.converter.convert(html, self.output_path)
        assert result.success, f"Should not crash on missing image: {result.error}"

    def test_img_without_src_does_not_crash(self):
        html = (
            "<section class='slide' data-type='content'>"
            "<h1>Empty Image</h1>"
            "<img alt='No Source'>"
            "</section>"
        )
        result = self.converter.convert(html, self.output_path)
        assert result.success, f"Should not crash on img without src: {result.error}"


class TestSlideElementParserImage:
    def test_img_tag_produces_image_in_slide_dict(self):
        from src.converters.base_parser import SlideElementParser
        parser = SlideElementParser()
        parser.feed("<h1>Test</h1><img src='chart.png' alt='Chart'>")
        slides = parser.get_slides()
        assert len(slides) >= 1
        slide = slides[0]
        images = slide.get("images", [])
        assert len(images) == 1
        assert images[0]["src"] == "chart.png"
        assert images[0]["alt"] == "Chart"

    def test_section_with_image(self):
        from src.converters.base_parser import SlideElementParser
        parser = SlideElementParser()
        parser.feed(
            "<section class='slide' data-type='content'>"
            "<h2>My Slide</h2>"
            "<img src='chart.png' alt='Chart'>"
            "</section>"
        )
        slides = parser.get_slides()
        assert len(slides) == 1
        assert slides[0]["slide_type"] == "content"
        assert slides[0]["title"] == "My Slide"
        assert len(slides[0]["images"]) == 1

    def test_multiple_slides_with_images(self):
        from src.converters.base_parser import SlideElementParser
        parser = SlideElementParser()
        parser.feed(
            "<section class='slide' data-type='content'>"
            "<h2>Slide 1</h2>"
            "<img src='a.png' alt='A'>"
            "</section>"
            "<section class='slide' data-type='data'>"
            "<h2>Slide 2</h2>"
            "<img src='b.png' alt='B'>"
            "<img src='c.png' alt='C'>"
            "</section>"
        )
        slides = parser.get_slides()
        assert len(slides) == 2
        assert len(slides[0]["images"]) == 1
        assert len(slides[1]["images"]) == 2
        assert slides[1]["slide_type"] == "data"


if __name__ == "__main__":
    import pytest
    pytest.main([__file__, "-v"])
