# -*- coding: utf-8 -*-
"""
PPT Path Verification Test
==========================

Verify that when output_format='pptx' is specified, the system actually
takes the PPT path through the entire pipeline:

  research_result → ContentOrchestrator.transform_to_html(output_format="pptx")
    → HTMLToPPTConverter.convert()
    → .pptx file with slides

  research_result → DocumentGenerator(config=format=PPTX)
    → _generate_pptx()
    → ContentOrchestrator.transform_to_html(output_format="pptx")
    → HTMLToPPTConverter.convert()
    → .pptx file with slides
"""
import sys
sys.path.insert(0, "E:/market_report_systerm")

import os
import tempfile
import shutil
import pytest
from pathlib import Path
from unittest.mock import patch, MagicMock


@pytest.fixture
def temp_dir():
    """Use a temp dir under cwd so path safety checks pass."""
    d = os.path.join(os.getcwd(), "__ppt_test_tmp__")
    os.makedirs(d, exist_ok=True)
    yield d
    shutil.rmtree(d, ignore_errors=True)


@pytest.fixture
def sample_research_result():
    return {
        "task_id": "test_pptx_001",
        "topic": "PPT Path Test",
        "title": "PPT Path Verification Report",
        "sections": [
            {
                "id": "market_overview",
                "title": "Market Overview",
                "content": "The market reached 950 billion in 2024, growing 37.5% year-over-year.\n\n## Key Data\n\n- Revenue: 950B\n- Growth: 37.5%",
                "order": 1,
                "type": "body",
                "data_points": [
                    {"metric": "Market Size", "value": "950", "unit": "B"},
                    {"metric": "Growth Rate", "value": "37.5", "unit": "%"},
                ],
            },
            {
                "id": "competition",
                "title": "Competitive Landscape",
                "content": "Top 3 players hold 65% market share.\n\n## Key Players\n\n1. Company A: 30%\n2. Company B: 20%\n3. Company C: 15%",
                "order": 2,
                "type": "body",
            },
        ],
        "key_findings": ["Market growing rapidly", "High concentration"],
    }


# ============================================================
# Test 1: ContentOrchestrator produces PPT-style HTML
# ============================================================
class TestContentOrchestratorPptxPath:
    """Verify ContentOrchestrator.transform_to_html with output_format='pptx'
    produces slide-based HTML, not article-based Word HTML.
    """

    def test_pptx_html_has_slides(self, sample_research_result):
        """PPT HTML should contain <section class='slide'> elements."""
        from src.content.content_orchestrator import ContentOrchestrator

        orchestrator = ContentOrchestrator()
        html_content = orchestrator.transform_to_html(
            research_result=sample_research_result,
            output_format="pptx",
        )

        assert html_content, "HTML content should not be empty"
        # PPT HTML uses slide-based sections, Word HTML uses article-based
        assert "slide" in html_content.lower(), (
            "PPT HTML should contain 'slide' class. "
            f"Got first 200 chars: {html_content[:200]}"
        )

    def test_pptx_html_differs_from_docx_html(self, sample_research_result):
        """PPT HTML and DOCX HTML should be structurally different."""
        from src.content.content_orchestrator import ContentOrchestrator

        orchestrator = ContentOrchestrator()
        pptx_html = orchestrator.transform_to_html(
            research_result=sample_research_result,
            output_format="pptx",
        )
        docx_html = orchestrator.transform_to_html(
            research_result=sample_research_result,
            output_format="docx",
        )

        assert pptx_html != docx_html, (
            "PPT and DOCX HTML should be structurally different. "
            "If they're the same, the format is being ignored."
        )


# ============================================================
# Test 2: HTMLToPPTConverter produces real .pptx with slides
# ============================================================
class TestHTMLToPPTConverterProducesPptx:
    """Verify HTMLToPPTConverter converts PPT HTML to .pptx with slides."""

    def test_converter_produces_pptx_file(self, sample_research_result, temp_dir):
        from src.content.content_orchestrator import ContentOrchestrator
        from src.converters.html_to_ppt import HTMLToPPTConverter

        # Step 1: Generate PPT HTML
        orchestrator = ContentOrchestrator()
        html_content = orchestrator.transform_to_html(
            research_result=sample_research_result,
            output_format="pptx",
        )

        # Step 2: Convert to PPTX
        converter = HTMLToPPTConverter()
        output_path = os.path.join(temp_dir, "test_report.pptx")
        result = converter.convert(html=html_content, output_path=output_path)

        assert result.success, f"PPTX conversion failed: {result.error}"
        assert os.path.exists(output_path), "PPTX file should exist"
        assert os.path.getsize(output_path) > 1000, "PPTX file should have content"

        # Step 3: Verify it's a real PPTX with slides
        from pptx import Presentation
        prs = Presentation(output_path)
        assert len(prs.slides) > 0, "PPTX should have at least 1 slide"

    def test_pptx_has_multiple_slides(self, sample_research_result, temp_dir):
        """PPTX should have multiple slides for multiple sections."""
        from src.content.content_orchestrator import ContentOrchestrator
        from src.converters.html_to_ppt import HTMLToPPTConverter

        orchestrator = ContentOrchestrator()
        html_content = orchestrator.transform_to_html(
            research_result=sample_research_result,
            output_format="pptx",
        )

        converter = HTMLToPPTConverter()
        output_path = os.path.join(temp_dir, "multi_slide.pptx")
        result = converter.convert(html=html_content, output_path=output_path)

        assert result.success

        from pptx import Presentation
        prs = Presentation(output_path)
        # Should have at least: cover + 2 content sections = 3 slides
        assert len(prs.slides) >= 2, (
            f"Expected >= 2 slides, got {len(prs.slides)}"
        )


# ============================================================
# Test 3: DocumentGenerator._generate_pptx uses HTMLToPPTConverter
# ============================================================
class TestDocumentGeneratorPptxPath:
    """Verify DocumentGenerator with format=PPTX calls _ppt_converter.convert(),
    NOT _word_converter.convert().
    """

    def test_generate_pptx_calls_ppt_converter(self, temp_dir):
        """DocumentGenerator.generate() with PPTX format should call _ppt_converter."""
        from src.core.orchestrator.output.document_generator import (
            DocumentGenerator,
            DocumentConfig,
            DocumentFormat,
        )
        from src.core.orchestrator.output.document_generator import DocumentResult

        config = DocumentConfig(format=DocumentFormat.PPTX, title="Test PPT")
        generator = DocumentGenerator(config)

        # Add minimal content
        generator.add_heading("Test Section", level=1)
        generator.add_paragraph("Test content")

        # Mock the PPT converter to track if it's called
        mock_ppt_result = MagicMock()
        mock_ppt_result.success = True
        mock_ppt_result.error = None

        with patch.object(generator, '_ppt_converter') as mock_ppt_conv, \
             patch.object(generator, '_word_converter') as mock_word_conv:

            mock_ppt_conv.convert.return_value = mock_ppt_result
            mock_word_conv.convert.return_value = MagicMock(success=False, error="should not be called")

            output_path = Path(temp_dir) / "test.pptx"
            result = generator.generate(output_path)

            # PPT converter MUST be called
            assert mock_ppt_conv.convert.called, (
                "DocumentGenerator with PPTX format must call _ppt_converter.convert()"
            )
            # Word converter MUST NOT be called
            assert not mock_word_conv.convert.called, (
                "DocumentGenerator with PPTX format must NOT call _word_converter.convert()"
            )

    def test_generate_docx_calls_word_converter(self, temp_dir):
        """DocumentGenerator.generate() with DOCX format should call _word_converter."""
        from src.core.orchestrator.output.document_generator import (
            DocumentGenerator,
            DocumentConfig,
            DocumentFormat,
        )

        config = DocumentConfig(format=DocumentFormat.DOCX, title="Test DOCX")
        generator = DocumentGenerator(config)
        generator.add_heading("Test Section", level=1)
        generator.add_paragraph("Test content")

        mock_word_result = MagicMock()
        mock_word_result.success = True
        mock_word_result.error = None

        with patch.object(generator, '_ppt_converter') as mock_ppt_conv, \
             patch.object(generator, '_word_converter') as mock_word_conv:

            mock_word_conv.convert.return_value = mock_word_result
            mock_ppt_conv.convert.return_value = MagicMock(success=False, error="should not be called")

            output_path = Path(temp_dir) / "test.docx"
            result = generator.generate(output_path)

            assert mock_word_conv.convert.called, "DOCX should call word converter"
            assert not mock_ppt_conv.convert.called, "DOCX should NOT call ppt converter"

    def test_generate_pptx_produces_real_file(self, sample_research_result, temp_dir):
        """End-to-end: DocumentGenerator PPTX produces a real .pptx file."""
        from src.core.orchestrator.output.document_generator import (
            DocumentGenerator,
            DocumentConfig,
            DocumentFormat,
        )

        config = DocumentConfig(format=DocumentFormat.PPTX, title="Test PPT")
        generator = DocumentGenerator(config)

        # Populate content
        for section in sample_research_result["sections"]:
            generator.add_heading(section["title"], level=1)
            generator.add_paragraph(section["content"])

        output_path = Path(temp_dir) / "e2e_test.pptx"
        result = generator.generate(output_path)

        assert result.path is not None, "Result path should not be None"
        assert output_path.exists(), f"PPTX file should exist at {output_path}"
        assert output_path.stat().st_size > 1000, "PPTX should have content"

        from pptx import Presentation
        prs = Presentation(str(output_path))
        assert len(prs.slides) > 0, "PPTX should have slides"


# ============================================================
# Test 4: DocumentGenerationAgent routes pptx to DocumentGenerator with PPTX format
# ============================================================
class TestDocumentGenerationAgentPptxRouting:
    """Verify DocumentGenerationAgent._handle_produce_document routes
    output_format='pptx' to DocumentGenerator with DocumentFormat.PPTX.
    """

    @pytest.mark.asyncio
    async def test_pptx_request_uses_pptx_format(self, sample_research_result, temp_dir):
        from src.agents.fixed_agents.document_generation_agent import DocumentGenerationAgent
        from src.core.orchestrator.output.document_generator import DocumentFormat as DGFormat

        agent = DocumentGenerationAgent(
            agent_id="test_agent",
            storage_path=temp_dir,
        )

        # Patch DocumentGenerator to capture the format it's called with
        captured_format = {}

        original_init = None
        from src.core.orchestrator.output import document_generator as dg_module

        class CapturingGenerator:
            def __init__(self, config):
                self.config = config
                captured_format["format"] = config.format
                self._content = []
                self._content_orchestrator = MagicMock()
                self._ppt_converter = MagicMock()
                self._word_converter = MagicMock()
                self._original_sections_meta = {}

            def add_heading(self, text, level=1):
                self._content.append({"type": "heading", "text": text, "level": level})

            def add_paragraph(self, text):
                self._content.append({"type": "paragraph", "text": text})

            def add_table(self, *args, **kwargs):
                pass

            def set_sections_meta(self, meta):
                pass

            def add_smart_chart(self, *args, **kwargs):
                return []

            def generate(self, output_path):
                from src.core.orchestrator.output.document_generator import DocumentResult
                Path(str(output_path)).write_bytes(b"fake pptx content")
                return DocumentResult(path=Path(str(output_path)), format=self.config.format)

        with patch.object(dg_module, "DocumentGenerator", CapturingGenerator):
            result = await agent.execute({
                "action": "produce_document",
                "research_result": sample_research_result,
                "output_format": "pptx",
                "output_dir": temp_dir,
                "task_id": "test_pptx_agent",
            })

        assert result.get("success"), f"Agent should succeed: {result}"
        assert captured_format.get("format") == DGFormat.PPTX, (
            f"DocumentGenerator should be called with PPTX format, "
            f"got {captured_format.get('format')}"
        )

    @pytest.mark.asyncio
    async def test_docx_request_uses_docx_format(self, sample_research_result, temp_dir):
        from src.agents.fixed_agents.document_generation_agent import DocumentGenerationAgent
        from src.core.orchestrator.output.document_generator import DocumentFormat as DGFormat

        agent = DocumentGenerationAgent(
            agent_id="test_agent",
            storage_path=temp_dir,
        )

        captured_format = {}
        from src.core.orchestrator.output import document_generator as dg_module

        class CapturingGenerator:
            def __init__(self, config):
                self.config = config
                captured_format["format"] = config.format
                self._content = []
                self._content_orchestrator = MagicMock()
                self._ppt_converter = MagicMock()
                self._word_converter = MagicMock()
                self._original_sections_meta = {}

            def add_heading(self, text, level=1):
                self._content.append({"type": "heading", "text": text, "level": level})

            def add_paragraph(self, text):
                self._content.append({"type": "paragraph", "text": text})

            def add_table(self, *args, **kwargs):
                pass

            def set_sections_meta(self, meta):
                pass

            def add_smart_chart(self, *args, **kwargs):
                return []

            def generate(self, output_path):
                from src.core.orchestrator.output.document_generator import DocumentResult
                Path(str(output_path)).write_bytes(b"fake docx content")
                return DocumentResult(path=Path(str(output_path)), format=self.config.format)

        with patch.object(dg_module, "DocumentGenerator", CapturingGenerator):
            result = await agent.execute({
                "action": "produce_document",
                "research_result": sample_research_result,
                "output_format": "docx",
                "output_dir": temp_dir,
                "task_id": "test_docx_agent",
            })

        assert result.get("success"), f"Agent should succeed: {result}"
        assert captured_format.get("format") == DGFormat.DOCX, (
            f"DocumentGenerator should be called with DOCX format, "
            f"got {captured_format.get('format')}"
        )


if __name__ == "__main__":
    import pytest
    pytest.main([__file__, "-v", "--tb=long"])
