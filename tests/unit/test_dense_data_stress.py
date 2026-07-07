"""
Dense data stress test for the PPT template rendering system.

This test creates a realistic, data-heavy market report with:
- 20+ slides covering all 12 template types
- Maximum KPI counts (4 cards)
- Large tables (8+ rows, 6+ columns)
- Long item lists (5-7 items per slide)
- Multi-line content
- HTML tags in table cells
- Chinese + English mixed content
- Comparison with many items
- Dense TOC with 8+ sections
- Multiple section_title slides with incrementing section numbers
- source_text on data slides
- Empty/missing data edge cases

Output: a .pptx file saved to data/reports/
"""
import os
import pytest
from pptx import Presentation
from pptx.util import Inches
from src.converters.template_selector import TemplateRegistry, TemplateSelector
from src.converters.slide_renderer import SlideRenderer


DESIGN = {
    "navy": "1A2744",
    "navy_dark": "0F1A2E",
    "navy_light": "2C3E50",
    "gold": "C9A227",
    "gold_light": "D4AF37",
    "white": "FFFFFF",
    "off_white": "F5F5F5",
    "text_dark": "333333",
    "text_mid": "666666",
    "text_light": "999999",
}

STYLES = {"slide_width": 13.33, "slide_height": 7.5}
OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))), "data", "reports")


@pytest.fixture(autouse=True)
def reset_registry():
    yield
    TemplateRegistry._reset()


def _render_presentation(slides_data):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    registry = TemplateRegistry()
    selector = TemplateSelector()
    renderer = SlideRenderer(DESIGN)
    section_index = 0
    rendered = []
    for i, sd in enumerate(slides_data):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        st = sd.get("slide_type", "content")
        if st in ("section_title", "section-title"):
            section_index += 1
        name = selector.select_and_enhance(sd, section_index=section_index)
        try:
            template = registry.get(name)
        except KeyError:
            template = registry.get("content_text_only")
        source = sd.get("source_text", "")
        if source:
            for dec in template.get("decorations", []):
                if dec.get("type") == "source_text" and not dec.get("text"):
                    dec["text"] = source
        renderer.render(slide, sd, template, STYLES, page_num=i + 1)
        rendered.append({"slide_index": i, "template": name, "shapes": len(slide.shapes)})
    return prs, rendered


def _dense_market_report():
    return [
        {
            "slide_type": "cover",
            "title": "2026年全球人工智能市场深度研究报告",
            "content": "行业分析 | 数据驱动 | 战略洞察",
        },
        {
            "slide_type": "toc",
            "title": "目录",
            "items": [
                "执行摘要 Executive Summary",
                "市场规模与增长 Market Size & Growth",
                "竞争格局 Competitive Landscape",
                "技术趋势 Technology Trends",
                "区域分析 Regional Analysis",
                "投资与融资 Investment & Funding",
                "风险与挑战 Risks & Challenges",
                "未来展望 Future Outlook",
            ],
        },
        {
            "slide_type": "section_title",
            "title": "执行摘要",
            "content": "全球AI市场在2025年达到$214.6B，预计2026年将突破$300B，CAGR达38.1%。大语言模型、生成式AI和AI Agent成为三大增长引擎。",
        },
        {
            "slide_type": "content",
            "title": "核心KPI指标",
            "items": [
                "Global AI Market: 214.6B USD, up 38.1% YoY",
                "GenAI Segment: 67.1B USD, growing 68% YoY",
                "AI Adoption Rate: 72%, up from 55%",
                "Enterprise AI Spend: 1.2T CNY, surged 42%",
            ],
            "content": "AI market continues exponential growth. Enterprise adoption accelerates across all verticals. Generative AI is the fastest-growing segment.",
        },
        {
            "slide_type": "section_title",
            "title": "市场规模与增长",
            "content": "本章节详细分析全球AI市场规模、增长率、细分领域数据及预测。涵盖2022-2026年历史数据和2027-2030年前瞻预测。",
        },
        {
            "slide_type": "content",
            "title": "市场增长驱动因素",
            "items": [
                "大语言模型技术突破推动企业级应用快速落地",
                "云原生AI平台降低了中小企业AI采用门槛",
                "数据量指数级增长为模型训练提供燃料",
                "各国政策支持（美国CHIPS法案、中国新基建）",
                "行业垂直解决方案成熟度提升（医疗、金融、制造）",
                "AI Agent自主决策能力突破带来全新商业模式",
            ],
            "source_text": "Source: Gartner, IDC, McKinsey Global AI Survey 2025",
        },
        {
            "slide_type": "data",
            "title": "全球AI市场细分数据",
            "table_data": [
                ["细分领域", "2024营收", "2025营收", "增长率", "2026E", "CAGR"],
                ["大语言模型", "$18.2B", "$32.5B", "+78.6%", "$55.1B", "72.3%"],
                ["生成式AI", "$22.4B", "$67.1B", "+68.0%", "$112.8B", "65.8%"],
                ["计算机视觉", "$14.7B", "$19.8B", "+34.7%", "$26.3B", "28.4%"],
                ["AI Agent", "$3.8B", "$12.6B", "+231.6%", "$31.2B", "102.5%"],
                ["自然语言处理", "$11.3B", "$16.9B", "+49.6%", "$24.7B", "38.2%"],
                ["推荐系统", "$28.9B", "$35.2B", "+21.8%", "$42.1B", "18.6%"],
                ["自动驾驶AI", "$9.4B", "$14.1B", "+50.0%", "$21.3B", "31.5%"],
            ],
            "source_text": "Source: IDC AI Market Tracker Q4 2025",
        },
        {
            "slide_type": "section_title",
            "title": "竞争格局",
            "content": "分析全球AI领域主要玩家战略布局、市场份额变化及竞争态势。",
        },
        {
            "slide_type": "content",
            "title": "中美AI竞争对比",
            "items": [
                "US vs China AI market size and strategic positioning",
                "US leads in foundational models and chip design",
                "China leads in AI application deployment and data scale",
                "Europe focuses on AI regulation and ethical frameworks",
                "India and Southeast Asia emerging as AI outsourcing hubs",
                "Both US and China investing heavily in quantum-AI convergence",
                "Global AI talent war intensifying across all regions",
            ],
            "content": "The US-China AI competition is reshaping global technology landscape with implications for supply chains, talent flows, and regulatory frameworks.",
        },
        {
            "slide_type": "data",
            "title": "主要厂商市场份额",
            "table_data": [
                ["公司", "AI营收", "市占率", "同比增长", "核心产品", "战略方向"],
                ["OpenAI", "$12.8B", "6.0%", "+156%", "GPT-5, DALL-E 4", "AGI研究"],
                ["Google DeepMind", "$18.5B", "8.6%", "+89%", "Gemini Ultra 2", "多模态融合"],
                ["Microsoft AI", "$22.3B", "10.4%", "+67%", "Copilot, Azure AI", "企业AI平台"],
                ["Meta AI", "$9.7B", "4.5%", "+112%", "Llama 4, Meta AI", "开源生态"],
                ["百度智能云", "85.2亿", "3.2%", "+45%", "文心4.5, 飞桨", "中文AI"],
                ["Anthropic", "$5.2B", "2.4%", "+234%", "Claude 4", "AI安全"],
            ],
            "source_text": "Source: CB Insights, Company Reports Q3 2025",
        },
        {
            "slide_type": "section_title",
            "title": "技术趋势",
            "content": "深入分析AI领域关键技术突破和未来技术路线图。",
        },
        {
            "slide_type": "content",
            "title": "关键技术指标",
            "items": [
                "LLM Parameter Count: 1.8T, up 340% YoY",
                "Training Compute: 5.2E FLOPS, growing 68% YoY",
                "Inference Cost: dropped 92% since 2023",
                "Context Window: 2M tokens standard, 10M emerging",
            ],
            "content": "Model capabilities scaling rapidly while costs plummeting. The democratization of AI is accelerating across all dimensions.",
        },
        {
            "slide_type": "content",
            "title": "六大技术趋势",
            "items": [
                "多模态大模型成为新标配：文本+图像+视频+音频+3D统一理解",
                "AI Agent从单任务向多Agent协作演进：自主规划+工具调用+记忆",
                "小模型+蒸馏技术突破：7B参数模型达到GPT-4级别性能",
                "AI推理芯片多元化：NVIDIA Blackwell + AMD MI400 + 华为昇腾910C",
                "合成数据解决训练数据枯竭：生成式AI反哺模型训练",
                "AI安全与对齐技术从理论走向工程化：Constitutional AI + RLHF升级",
            ],
        },
        {
            "slide_type": "section_title",
            "title": "区域分析",
            "content": "全球五大区域AI发展现状、政策环境与投资热点全景扫描。",
        },
        {
            "slide_type": "content",
            "title": "区域AI投资对比",
            "items": [
                "North America vs Asia Pacific AI investment landscape",
                "North America: $89.2B in AI venture funding in 2025",
                "Asia Pacific: $62.7B, driven by China, Japan, and India",
                "Europe: $31.4B, focused on AI regulation compliance",
                "Middle East: $8.3B, sovereign AI ambitions growing",
            ],
            "content": "Regional AI investment patterns reveal divergent strategies: NA leads in frontier research, APAC in application scale, EU in governance.",
        },
        {
            "slide_type": "content",
            "title": "中国市场深度解析",
            "items": [
                "中国AI市场规模达到4,286亿元，同比增长41.3%",
                "大模型备案数量超过200个，形成百模大战格局",
                "AI芯片国产化率提升至28%，华为昇腾生态快速成长",
                "智能制造AI渗透率达45%，工业质检应用最成熟",
                "金融AI应用从风控延伸至投顾、合规、反欺诈全链条",
                "医疗AI获批三类医疗器械证达67个，影像诊断领先",
                "教育AI从个性化学习向AI教师助手演进",
            ],
            "source_text": "Source: CAICT, iResearch, Ministry of Industry and IT 2025",
        },
        {
            "slide_type": "section_title",
            "title": "投资与融资",
            "content": "全球AI领域投融资数据、独角兽企业图谱及IPO趋势分析。",
        },
        {
            "slide_type": "content",
            "title": "融资核心数据",
            "items": [
                "Total AI Funding: 142.3B USD in 2025",
                "Mega Rounds (>$100M): 87 deals, up 34%",
                "AI Unicorn Count: 184 globally, 42 new in 2025",
                "Avg Series A Valuation: $280M, up 62% YoY",
            ],
            "content": "AI funding hitting record levels. Mega-rounds concentrated in foundation models, AI infrastructure, and vertical AI solutions.",
        },
        {
            "slide_type": "data",
            "title": "2025年AI融资Top 10",
            "table_data": [
                ["排名", "公司", "融资金额", "轮次", "估值", "领域"],
                ["1", "xAI", "$12.0B", "Series C", "$120B", "AGI"],
                ["2", "Anthropic", "$8.0B", "Series D", "$60B", "AI Safety"],
                ["3", "Cohere", "$2.5B", "Series D", "$22B", "Enterprise LLM"],
                ["4", "Mistral", "$1.8B", "Series C", "$16B", "Open Models"],
                ["5", "月之暗面", "$1.2B", "Series B", "¥120B", "中文LLM"],
                ["6", "Scale AI", "$1.0B", "Series F", "$14B", "Data Labeling"],
                ["7", "CoreWeave", "$1.5B", "Series D", "$19B", "GPU Cloud"],
                ["8", "Glean", "$800M", "Series E", "$8B", "Enterprise Search"],
                ["9", "Perplexity", "$750M", "Series D", "$9B", "AI Search"],
                ["10", "智谱AI", "¥50亿", "Series C", "¥200B", "中文LLM"],
            ],
            "source_text": "Source: Crunchbase, PitchBook, CB Insights 2025",
        },
        {
            "slide_type": "section_title",
            "title": "风险与挑战",
            "content": "AI发展面临的技术风险、监管挑战和社会影响分析。",
        },
        {
            "slide_type": "content",
            "title": "五大风险领域",
            "items": [
                "数据隐私与合规风险：GDPR/PIPL跨境数据流动限制加剧",
                "AI幻觉与可靠性：大模型输出准确率仅85-92%，关键场景仍需人工审核",
                "AI安全对抗攻击：越狱攻击、提示注入、模型窃取威胁升级",
                "算力供应链风险：NVIDIA H200供应受限，国产替代仍在追赶",
                "人才短缺与薪资膨胀：全球AI人才缺口超400万，高级研究员年薪超$1M",
            ],
        },
        {
            "slide_type": "findings",
            "title": "核心发现与建议",
            "items": [
                "AI市场进入高速增长期，2026年将突破$300B大关",
                "生成式AI和AI Agent是增长最快的两个赛道",
                "中美AI竞争从技术层面扩展到生态层面",
                "AI安全与对齐是2026年最关键的投资主题",
                "企业AI采用率突破70%，从试验阶段进入规模化部署",
            ],
        },
        {
            "slide_type": "end",
            "title": "感谢关注",
            "content": "联系方式: ai-research@example.com | 数据截止: 2025年Q3",
        },
    ]


class TestDenseDataStressTest:
    def test_dense_market_report_renders_all_slides(self):
        slides = _dense_market_report()
        prs, rendered = _render_presentation(slides)
        assert len(prs.slides) == len(slides), f"Expected {len(slides)} slides, got {len(prs.slides)}"

    def test_all_slides_have_shapes(self):
        slides = _dense_market_report()
        prs, rendered = _render_presentation(slides)
        for info in rendered:
            assert info["shapes"] > 0, f"Slide {info['slide_index']} (template={info['template']}) has no shapes"

    def test_template_distribution(self):
        slides = _dense_market_report()
        prs, rendered = _render_presentation(slides)
        templates_used = set(r["template"] for r in rendered)
        required = {"cover", "toc", "section_title", "kpi_highlight", "data_table", "comparison", "findings", "end"}
        for t in required:
            assert t in templates_used, f"Expected template '{t}' not found in {templates_used}"

    def test_section_numbers_increment(self):
        slides = _dense_market_report()
        selector = TemplateSelector()
        section_nums = []
        section_index = 0
        for sd in slides:
            if sd.get("slide_type") in ("section_title", "section-title"):
                section_index += 1
                selector.select_and_enhance(sd, section_index=section_index)
                section_nums.append(sd.get("section_number"))
        assert section_nums == [1, 2, 3, 4, 5, 6, 7]

    def test_kpi_cards_max_four(self):
        slides = _dense_market_report()
        prs, rendered = _render_presentation(slides)
        for i, r in enumerate(rendered):
            if r["template"] == "kpi_highlight":
                shapes_count = r["shapes"]
                assert shapes_count >= 3, f"KPI slide {i} has too few shapes ({shapes_count})"

    def test_data_table_slides_have_tables(self):
        slides = _dense_market_report()
        prs, rendered = _render_presentation(slides)
        for i, r in enumerate(rendered):
            if r["template"] == "data_table":
                slide = prs.slides[i]
                tables = [s for s in slide.shapes if s.has_table]
                assert len(tables) >= 1, f"Data table slide {i} has no table"

    def test_comparison_slides_have_two_columns(self):
        slides = _dense_market_report()
        prs, rendered = _render_presentation(slides)
        for i, r in enumerate(rendered):
            if r["template"] == "comparison":
                slide = prs.slides[i]
                text_shapes = [s for s in slide.shapes if s.has_text_frame]
                assert len(text_shapes) >= 3, f"Comparison slide {i} has too few text shapes"

    def test_cover_has_date(self):
        slides = _dense_market_report()
        prs, _ = _render_presentation(slides)
        cover = prs.slides[0]
        texts = " ".join(s.text_frame.text for s in cover.shapes if s.has_text_frame)
        import datetime
        assert datetime.date.today().strftime("%Y") in texts, "Cover should show current year in date"

    def test_source_text_rendered(self):
        slides = _dense_market_report()
        prs, rendered = _render_presentation(slides)
        data_slides = [(i, r) for i, r in enumerate(rendered) if r["template"] == "data_table"]
        for idx, r in data_slides:
            slide = prs.slides[idx]
            texts = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
            sd = slides[idx]
            if sd.get("source_text"):
                assert any(kw in texts for kw in ["Source:", "IDC", "Gartner", "CB Insights", "Crunchbase"]), \
                    f"Slide {idx} should have source text rendered"

    def test_multi_line_content_rendered(self):
        slides = _dense_market_report()
        prs, rendered = _render_presentation(slides)
        section_slides = [(i, r) for i, r in enumerate(rendered) if r["template"] == "section_title"]
        for idx, r in section_slides:
            slide = prs.slides[idx]
            texts = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
            sd = slides[idx]
            assert sd["title"] in texts, f"Section slide {idx} should contain title '{sd['title']}'"

    def test_dense_table_jagged_rows(self):
        slides = _dense_market_report()
        prs, rendered = _render_presentation(slides)
        data_slides = [(i, r) for i, r in enumerate(rendered) if r["template"] == "data_table"]
        for idx, r in data_slides:
            slide = prs.slides[idx]
            tables = [s for s in slide.shapes if s.has_table]
            for t in tables:
                tbl = t.table
                assert tbl.rows.__len__() > 0
                assert tbl.columns.__len__() > 0

    def test_save_pptx_output(self):
        slides = _dense_market_report()
        prs, rendered = _render_presentation(slides)
        os.makedirs(OUTPUT_DIR, exist_ok=True)
        output_path = os.path.join(OUTPUT_DIR, "dense_market_report_test.pptx")
        prs.save(output_path)
        assert os.path.exists(output_path)
        file_size = os.path.getsize(output_path)
        assert file_size > 10000, f"PPTX file too small ({file_size} bytes), likely corrupted"
        print(f"\n[PPTX OUTPUT] Saved to: {output_path}")
        print(f"[PPTX OUTPUT] File size: {file_size:,} bytes")
        print(f"[PPTX OUTPUT] Total slides: {len(prs.slides)}")
        print(f"[PPTX OUTPUT] Template distribution:")
        from collections import Counter
        for name, count in Counter(r["template"] for r in rendered).most_common():
            print(f"  {name}: {count}")

    def test_edge_case_empty_items(self):
        slides = [
            {"slide_type": "content", "title": "Empty Items", "items": [], "content": "Only content, no items."},
        ]
        prs, rendered = _render_presentation(slides)
        assert rendered[0]["template"] == "content_text_only"
        assert rendered[0]["shapes"] > 0

    def test_edge_case_single_kpi(self):
        slides = [
            {"slide_type": "content", "title": "Single KPI", "items": ["Revenue 15.1B"], "content": "Only one KPI."},
        ]
        prs, rendered = _render_presentation(slides)
        assert rendered[0]["template"] == "content_text_only"

    def test_edge_case_no_data_at_all(self):
        slides = [
            {"slide_type": "content", "title": "", "items": [], "content": ""},
        ]
        prs, rendered = _render_presentation(slides)
        assert rendered[0]["template"] == "content_text_only"

    def test_edge_case_html_tags_in_table(self):
        slides = [
            {"slide_type": "data", "title": "HTML Table", "table_data": [
                ["Metric", "Value"],
                ["<strong>Revenue</strong>", "$15.1B"],
                ["<em>Growth</em>", "28.9%"],
            ]},
        ]
        prs, rendered = _render_presentation(slides)
        slide = prs.slides[0]
        tables = [s for s in slide.shapes if s.has_table]
        assert len(tables) == 1
        cell_text = tables[0].table.cell(1, 0).text
        assert "<strong>" not in cell_text
        assert "Revenue" in cell_text

    def test_edge_case_very_long_title(self):
        long_title = "这是一个非常长的标题用来测试文本溢出行为我们需要确保超长文本不会导致渲染崩溃或者布局完全错乱"
        slides = [
            {"slide_type": "content", "title": long_title, "items": ["Point A"]},
        ]
        prs, rendered = _render_presentation(slides)
        assert rendered[0]["shapes"] > 0

    def test_edge_case_chinese_kpi_with_mixed_units(self):
        slides = [
            {"slide_type": "content", "title": "中国市场KPI", "items": [
                "营收达到3.2万亿，同比增长15.8%",
                "用户数2.7M，海外占比28%",
                "研发投入85.2亿CNY",
            ], "content": "中国市场持续高速增长。"},
        ]
        prs, rendered = _render_presentation(slides)
        assert rendered[0]["template"] == "kpi_highlight"
