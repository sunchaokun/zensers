import os
import pytest
from pptx import Presentation
from pptx.util import Inches
from src.converters.template_selector import TemplateRegistry, TemplateSelector
from src.converters.slide_renderer import SlideRenderer


DESIGN = {
    "navy": "1A2744",
    "navy_dark": "0F1A2E",
    "navy_light": "2C3E50",
    "gold": "C9A227",
    "gold_light": "D4AF37",
    "white": "FFFFFF",
    "off_white": "F5F5F5",
    "text_dark": "333333",
    "text_mid": "666666",
    "text_light": "999999",
}

STYLES = {"slide_width": 13.33, "slide_height": 7.5}


@pytest.fixture(autouse=True)
def reset_registry():
    yield
    TemplateRegistry._reset()


def _render_slides(slides_data, section_start=0):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    registry = TemplateRegistry()
    selector = TemplateSelector()
    renderer = SlideRenderer(DESIGN)
    section_index = section_start
    for i, sd in enumerate(slides_data):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        st = sd.get("slide_type", "content")
        if st in ("section_title", "section-title"):
            section_index += 1
        name = selector.select_and_enhance(sd, section_index=section_index)
        template = registry.get(name)
        renderer.render(slide, sd, template, STYLES, page_num=i + 1)
    return prs


class TestIntegrationFullPresentation:
    def test_full_presentation_all_template_types(self):
        slides_data = [
            {"slide_type": "cover", "title": "Market Report 2026", "content": "Industry Analysis"},
            {"slide_type": "toc", "title": "Table of Contents", "items": ["Executive Summary", "Market Overview", "Competitive Landscape", "Future Outlook"]},
            {"slide_type": "section_title", "title": "Market Overview", "content": "This section covers market size, growth trends, and key dynamics shaping the industry."},
            {"slide_type": "content", "title": "Key Metrics", "items": ["Revenue 15.1B USD", "Users 2.7M", "Growth 28.9%"], "content": "Strong growth. Revenue exceeded expectations."},
            {"slide_type": "content", "title": "Market Drivers", "items": ["Digital transformation", "Cloud adoption", "AI integration"]},
            {"slide_type": "data", "title": "Financial Data", "table_data": [["Metric", "2024", "2025"], ["Revenue", "$12B", "$15.1B"], ["Growth", "22%", "28.9%"]]},
            {"slide_type": "content", "title": "US vs China", "items": ["US vs China market", "US leads in tech", "China leads in manufacturing", "Both growing rapidly", "India emerging"]},
            {"slide_type": "findings", "title": "Key Findings", "items": ["Market growing at 28.9% CAGR", "AI segment fastest growing", "Cloud migration accelerating"]},
            {"slide_type": "end", "title": "Thank You", "content": "Questions & Discussion"},
        ]
        prs = _render_slides(slides_data)
        assert len(prs.slides) == 9
        for slide in prs.slides:
            assert len(slide.shapes) > 0

    def test_kpi_slide_has_cards(self):
        slides_data = [
            {"slide_type": "content", "title": "KPIs", "items": ["Revenue 15.1B", "Users 2.7M", "Growth 28.9%"], "content": "Strong growth this year."},
        ]
        prs = _render_slides(slides_data)
        assert len(prs.slides) == 1
        texts = []
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        all_text = " ".join(texts)
        assert "15.1B" in all_text
        assert "2.7M" in all_text

    def test_section_title_zero_padded(self):
        slides_data = [
            {"slide_type": "section_title", "title": "Executive Summary", "content": "Overview of key findings."},
        ]
        prs = _render_slides(slides_data, section_start=0)
        texts = []
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        all_text = " ".join(texts)
        assert "01" in all_text
        assert "Executive Summary" in all_text

    def test_data_table_slide(self):
        slides_data = [
            {"slide_type": "data", "title": "Financials", "table_data": [["Year", "Revenue"], ["2024", "$12B"], ["2025", "$15.1B"]]},
        ]
        prs = _render_slides(slides_data)
        tables = [s for s in prs.slides[0].shapes if s.has_table]
        assert len(tables) == 1
        assert tables[0].table.cell(0, 0).text == "Year"

    def test_feature_flag_default_off(self):
        from src.converters.html_to_ppt import HTMLToPPTConverter
        converter = HTMLToPPTConverter()
        assert not converter._should_use_template_renderer()

    def test_feature_flag_env_on(self, monkeypatch):
        monkeypatch.setenv("USE_TEMPLATE_RENDERER", "1")
        from src.converters.html_to_ppt import HTMLToPPTConverter
        converter = HTMLToPPTConverter()
        assert converter._should_use_template_renderer()
