import json
import pytest
from pptx import Presentation
from pptx.util import Inches
from src.converters.slide_renderer import SlideRenderer


DESIGN = {
    "navy": "1A2744",
    "navy_dark": "0F1A2E",
    "navy_light": "2C3E50",
    "gold": "C9A227",
    "gold_light": "D4AF37",
    "white": "FFFFFF",
    "off_white": "F5F5F5",
    "text_dark": "333333",
    "text_mid": "666666",
    "text_light": "999999",
}


@pytest.fixture
def renderer():
    return SlideRenderer(DESIGN)


@pytest.fixture
def slide():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


class TestResolveColor:
    def test_named_color_from_design(self, renderer):
        assert renderer._resolve_color("navy") == "1A2744"

    def test_hex_color_passthrough(self, renderer):
        assert renderer._resolve_color("4CAF50") == "4CAF50"

    def test_unknown_named_color_passthrough(self, renderer):
        assert renderer._resolve_color("coral") == "coral"


class TestRenderBackground:
    def test_solid_background(self, renderer, slide):
        bg = {"type": "solid", "color": "navy"}
        renderer._render_background(slide, bg)
        fill = slide.background.fill
        assert fill.type is not None

    def test_gradient_background(self, renderer, slide):
        bg = {"type": "gradient", "color1": "navy", "color2": "navy_light"}
        renderer._render_background(slide, bg)
        fill = slide.background.fill
        assert fill.type is not None

    def test_empty_background_no_error(self, renderer, slide):
        renderer._render_background(slide, {})


class TestRender:
    def test_render_calls_all_layers(self, renderer, slide):
        template = {
            "background": {"type": "solid", "color": "white"},
            "slots": [
                {"id": "title", "type": "text", "source": "title", "position": {"left": 0.8, "top": 0.3, "width": 11.7, "height": 0.7}, "style": {"font_size": 24, "font_weight": "bold", "color": "navy"}}
            ],
            "decorations": [
                {"type": "footer_bar", "layer": "bottom", "color": "gold", "height": 0.11},
                {"type": "page_number", "layer": "top", "position": "bottom_right", "color": "text_light", "font_size": 10}
            ]
        }
        slide_data = {"title": "Test Title", "items": [], "images": []}
        styles = {"slide_width": 13.33, "slide_height": 7.5}
        renderer.render(slide, slide_data, template, styles, page_num=1)
        assert len(slide.shapes) > 0


class TestRenderTextSlot:
    def test_title_rendered(self, renderer, slide):
        slot = {"id": "title", "type": "text", "source": "title", "position": {"left": 0.8, "top": 0.3, "width": 11.7, "height": 0.7}, "style": {"font_size": 24, "font_weight": "bold", "color": "navy"}}
        renderer._render_text_slot(slide, slot, {"title": "Hello World"}, {})
        assert any("Hello World" in s.text_frame.text for s in slide.shapes if s.has_text_frame)

    def test_section_number_zero_padded(self, renderer, slide):
        slot = {"id": "sn", "type": "text", "source": "section_number", "position": {"left": 0.8, "top": 0.5, "width": 3, "height": 1}, "style": {"font_size": 48, "format": "zero_padded", "color": "gold"}}
        renderer._render_text_slot(slide, slot, {"section_number": 3}, {})
        assert any("03" in s.text_frame.text for s in slide.shapes if s.has_text_frame)

    def test_empty_source_skips(self, renderer, slide):
        slot = {"id": "t", "type": "text", "source": "title", "position": {"left": 0.8, "top": 0.3, "width": 11.7, "height": 0.7}, "style": {"font_size": 24}}
        renderer._render_text_slot(slide, slot, {"title": ""}, {})
        assert len(slide.shapes) == 0


class TestRenderItemsSlot:
    def test_items_rendered(self, renderer, slide):
        slot = {"id": "bi", "type": "items", "source": "items", "position": {"left": 0.8, "top": 1.3, "width": 11.7, "height": 5.2}, "style": {"bullet": "▸", "bullet_color": "gold", "font_size": 14, "color": "text_dark", "max_items": 5}}
        renderer._render_items_slot(slide, slot, {"items": ["Point A", "Point B"]}, {})
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("Point A" in t for t in texts)
        assert any("Point B" in t for t in texts)

    def test_max_items_respected(self, renderer, slide):
        slot = {"id": "bi", "type": "items", "source": "items", "position": {"left": 0.8, "top": 1.3, "width": 11.7, "height": 5.2}, "style": {"bullet": "▸", "max_items": 2, "font_size": 14}}
        renderer._render_items_slot(slide, slot, {"items": ["A", "B", "C"]}, {})
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert not any("C" in t for t in texts)

    def test_empty_items_skips(self, renderer, slide):
        slot = {"id": "bi", "type": "items", "source": "items", "position": {}, "style": {}}
        renderer._render_items_slot(slide, slot, {"items": []}, {})
        assert len(slide.shapes) == 0


class TestRenderTableSlot:
    def test_table_rendered(self, renderer, slide):
        slot = {"id": "dt", "type": "table", "source": "table_data", "position": {"left": 0.8, "top": 1.3, "width": 5.6, "height": "auto"}, "style": {"header_bg": "navy", "header_color": "white", "row_font_size": 12}}
        renderer._render_table_slot(slide, slot, {"table_data": [["H1", "H2"], ["A", "B"]]}, {})
        tables = [s for s in slide.shapes if s.has_table]
        assert len(tables) == 1
        assert tables[0].table.cell(0, 0).text == "H1"

    def test_empty_table_skips(self, renderer, slide):
        slot = {"id": "dt", "type": "table", "source": "table_data", "position": {}, "style": {}}
        renderer._render_table_slot(slide, slot, {"table_data": []}, {})
        assert len(slide.shapes) == 0

    def test_table_full_width_when_no_image(self, renderer, slide):
        slot = {"id": "dt", "type": "table", "source": "table_data", "position": {"left": 0.8, "top": 1.3, "width": 5.6, "height": "auto"}, "style": {"header_bg": "navy", "header_color": "white", "row_font_size": 12, "full_width_when_no_image": True}}
        slide_data = {"table_data": [["H1", "H2"], ["A", "B"]]}
        styles = {"slide_width": 13.33}
        renderer._render_table_slot(slide, slot, slide_data, styles)
        tables = [s for s in slide.shapes if s.has_table]
        assert len(tables) == 1
        expected_width = 13.33 - 0.8 - 0.8
        actual_width = tables[0].width / 914400
        assert abs(actual_width - expected_width) < 0.01

    def test_table_half_width_with_image(self, renderer, slide):
        slot = {"id": "dt", "type": "table", "source": "table_data", "position": {"left": 0.8, "top": 1.3, "width": 5.6, "height": "auto"}, "style": {"header_bg": "navy", "header_color": "white", "row_font_size": 12, "full_width_when_no_image": True}}
        slide_data = {"table_data": [["H1", "H2"], ["A", "B"]], "images": [{"src": "chart.png"}]}
        styles = {"slide_width": 13.33}
        renderer._render_table_slot(slide, slot, slide_data, styles)
        tables = [s for s in slide.shapes if s.has_table]
        assert len(tables) == 1
        actual_width = tables[0].width / 914400
        assert abs(actual_width - 5.6) < 0.01


class TestRenderKpiCardsSlot:
    def test_kpi_cards_rendered(self, renderer, slide):
        slot = {"id": "kr", "type": "kpi_cards", "source": "kpi_data", "position": {"left": 0.8, "top": 1.5, "width": 11.7, "height": 3.5}, "style": {"card_bg": "navy", "number_size": 36, "number_color": "gold", "label_size": 12, "label_color": "white", "max_cards": 4}}
        kpi_data = [{"number": "15.1B", "label": "Revenue", "trend": None, "trend_direction": None}]
        renderer._render_kpi_cards_slot(slide, slot, {"kpi_data": kpi_data}, {})
        shapes = [s for s in slide.shapes if s.has_text_frame]
        assert len(shapes) >= 1
        assert any("15.1B" in s.text_frame.text for s in shapes)

    def test_kpi_with_trend(self, renderer, slide):
        slot = {"id": "kr", "type": "kpi_cards", "source": "kpi_data", "position": {"left": 0.8, "top": 1.5, "width": 11.7, "height": 3.5}, "style": {"card_bg": "navy", "number_size": 36, "number_color": "gold", "label_size": 12, "label_color": "white", "trend_up": "↑", "trend_down": "↓", "max_cards": 4}}
        kpi_data = [{"number": "15.1B", "label": "Revenue", "trend": "28.9%", "trend_direction": "up"}]
        renderer._render_kpi_cards_slot(slide, slot, {"kpi_data": kpi_data}, {})
        texts = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
        assert "15.1B" in texts
        assert "↑" in texts


class TestRenderInsightBarSlot:
    def test_insight_bar_rendered(self, renderer, slide):
        slot = {"id": "ins", "type": "insight_bar", "source": "insight_text", "position": {"left": 0.8, "top": 5.8, "width": 11.7, "height": 1.0}, "style": {"bg_color": "navy", "icon": "💡", "font_size": 13, "color": "white"}}
        renderer._render_insight_bar_slot(slide, slot, {"insight_text": "Key insight here"}, {})
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("Key insight here" in t for t in texts)

    def test_no_insight_skips(self, renderer, slide):
        slot = {"id": "ins", "type": "insight_bar", "source": "insight_text", "position": {}, "style": {}}
        renderer._render_insight_bar_slot(slide, slot, {"insight_text": ""}, {})
        assert len(slide.shapes) == 0


class TestRenderComparisonSlot:
    def test_comparison_rendered(self, renderer, slide):
        slot = {"id": "cmp", "type": "comparison", "source": "comparison_data", "position": {"left": 0.8, "top": 1.3, "width": 11.7, "height": 5.2}, "style": {"left_color": "navy", "right_color": "gold", "font_size": 13}}
        comp_data = {"left": {"title": "US", "items": ["Strong"]}, "right": {"title": "China", "items": ["Growing"]}}
        renderer._render_comparison_slot(slide, slot, {"comparison_data": comp_data}, {})
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("US" in t for t in texts)
        assert any("China" in t for t in texts)

    def test_no_comparison_skips(self, renderer, slide):
        slot = {"id": "cmp", "type": "comparison", "source": "comparison_data", "position": {}, "style": {}}
        renderer._render_comparison_slot(slide, slot, {"comparison_data": {}}, {})
        assert len(slide.shapes) == 0


class TestRenderDecoration:
    def test_footer_bar(self, renderer, slide):
        dec = {"type": "footer_bar", "layer": "bottom", "color": "gold", "height": 0.11}
        renderer._render_decoration(slide, dec, {"slide_width": 13.33, "slide_height": 7.5})
        assert len(slide.shapes) == 1

    def test_side_accent(self, renderer, slide):
        dec = {"type": "side_accent", "layer": "bottom", "color": "gold", "width": 0.06}
        renderer._render_decoration(slide, dec, {"slide_width": 13.33, "slide_height": 7.5})
        assert len(slide.shapes) == 1

    def test_title_underline(self, renderer, slide):
        dec = {"type": "title_underline", "layer": "bottom", "color": "gold", "width": 4.0, "offset_top": 1.05}
        renderer._render_decoration(slide, dec, {"slide_width": 13.33, "slide_height": 7.5})
        assert len(slide.shapes) == 1

    def test_page_number(self, renderer, slide):
        dec = {"type": "page_number", "layer": "top", "position": "bottom_right", "color": "text_light", "font_size": 10}
        renderer._render_decoration(slide, dec, {"slide_width": 13.33, "slide_height": 7.5}, page_num=3)
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("3" in t for t in texts)

    def test_page_number_zero_skips(self, renderer, slide):
        dec = {"type": "page_number", "layer": "top", "position": "bottom_right", "color": "text_light", "font_size": 10}
        renderer._render_decoration(slide, dec, {"slide_width": 13.33, "slide_height": 7.5}, page_num=0)
        assert len(slide.shapes) == 0

    def test_branding(self, renderer, slide):
        dec = {"type": "branding", "layer": "top", "text": "CONFIDENTIAL", "position": "top_right", "color": "text_light", "font_size": 9}
        renderer._render_decoration(slide, dec, {"slide_width": 13.33, "slide_height": 7.5})
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("CONFIDENTIAL" in t for t in texts)

    def test_source_text_empty_skips(self, renderer, slide):
        dec = {"type": "source_text", "layer": "top", "text": "", "position": "bottom_left", "color": "text_light", "font_size": 9}
        renderer._render_decoration(slide, dec, {"slide_width": 13.33, "slide_height": 7.5})
        assert len(slide.shapes) == 0
