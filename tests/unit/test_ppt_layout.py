import pytest
import os
import tempfile
from pptx import Presentation
from pptx.util import Inches
from src.converters.html_to_ppt import HTMLToPPTConverter


@pytest.fixture
def converter():
    return HTMLToPPTConverter()


@pytest.fixture
def styles():
    return {"slide_width": 13.333, "slide_height": 7.5}


def _make_pptx_with_slide(converter, slide_data, styles):
    prs = Presentation()
    prs.slide_width = Inches(styles["slide_width"])
    prs.slide_height = Inches(styles["slide_height"])
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    converter._create_content_slide(slide, slide_data, styles)
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        prs.save(f.name)
        return f.name, prs


class TestLeftTextRightImageLayout:
    def test_text_box_on_left_half(self, converter, styles):
        slide_data = {
            "slide_type": "content",
            "title": "Market Overview",
            "items": ["Point 1", "Point 2", "Point 3"],
            "images": [{"src": "nonexistent.png", "alt": ""}],
        }
        filepath, prs = _make_pptx_with_slide(converter, slide_data, styles)
        try:
            slide = prs.slides[0]
            text_boxes = [s for s in slide.shapes if hasattr(s, "text_frame") and s.text_frame.text.strip()]
            assert len(text_boxes) >= 2
            content_box = None
            for tb in text_boxes:
                if "Point" in tb.text_frame.text:
                    content_box = tb
                    break
            assert content_box is not None
            left_inches = content_box.left / 914400
            width_inches = content_box.width / 914400
            assert left_inches < 1.5
            assert width_inches < 7.0
        finally:
            os.unlink(filepath)

    def test_no_image_text_uses_full_width(self, converter, styles):
        slide_data = {
            "slide_type": "content",
            "title": "Summary",
            "items": ["Point 1", "Point 2"],
        }
        filepath, prs = _make_pptx_with_slide(converter, slide_data, styles)
        try:
            slide = prs.slides[0]
            text_boxes = [s for s in slide.shapes if hasattr(s, "text_frame") and "Point" in s.text_frame.text]
            assert len(text_boxes) == 1
            width_inches = text_boxes[0].width / 914400
            assert width_inches > 10.0
        finally:
            os.unlink(filepath)

    def test_findings_slide_left_text_right_image(self, converter, styles):
        prs = Presentation()
        prs.slide_width = Inches(styles["slide_width"])
        prs.slide_height = Inches(styles["slide_height"])
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)
        slide_data = {
            "slide_type": "findings",
            "title": "Key Findings",
            "items": ["Finding 1", "Finding 2", "Finding 3"],
            "images": [{"src": "nonexistent.png", "alt": ""}],
        }
        converter._create_findings_slide(slide, slide_data, styles)
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs.save(f.name)
            filepath = f.name
        try:
            slide2 = Presentation(filepath).slides[0]
            text_boxes = [s for s in slide2.shapes if hasattr(s, "text_frame") and "Finding" in s.text_frame.text and "Key Findings" not in s.text_frame.text]
            assert len(text_boxes) == 1
            width_inches = text_boxes[0].width / 914400
            assert width_inches < 7.0
        finally:
            os.unlink(filepath)

    def test_data_slide_left_table_right_image(self, converter, styles):
        prs = Presentation()
        prs.slide_width = Inches(styles["slide_width"])
        prs.slide_height = Inches(styles["slide_height"])
        blank = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank)
        slide_data = {
            "slide_type": "data",
            "title": "Data Table",
            "table_data": [["Metric", "Value"], ["Revenue", "$1.5B"]],
            "images": [{"src": "nonexistent.png", "alt": ""}],
        }
        converter._create_data_slide(slide, slide_data, styles)
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs.save(f.name)
            filepath = f.name
        try:
            slide2 = Presentation(filepath).slides[0]
            tables = [s for s in slide2.shapes if s.has_table]
            assert len(tables) == 1
            width_inches = tables[0].width / 914400
            assert width_inches < 7.0
        finally:
            os.unlink(filepath)
