"""Generate a real PPT using the data-driven pipeline.

Usage: python scripts/generate_ppt_from_data.py
"""
import sys
import os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pathlib import Path
from src.core.adjustment.ppt_input_adapter import PptInputAdapter
from src.core.adjustment.ppt_requirement_extractor import PptRequirementExtractor
from src.core.adjustment.ppt_data_supplementer import PptDataSupplementer
from src.core.adjustment.slide_data_builder import SlideDataBuilder
from src.core.adjustment.ppt_structure_editor import PptStructureEditor
from src.core.adjustment.slide_data_store import SlideDataStore
from src.core.adjustment.ppt_version_manager import PptVersionManager
from src.content.content_orchestrator import ContentSection, SectionType


def create_sample_docx():
    from docx import Document

    doc = Document()
    doc.add_heading("2026年中国新能源汽车产业深度研究报告", level=0)
    doc.add_paragraph("产业链全景分析 · 竞争格局研判 · 投资价值评估")

    doc.add_heading("一、行业概览", level=1)
    doc.add_paragraph(
        "2025年中国新能源汽车销量达到950万辆，同比增长37.5%，渗透率突破40%。"
        "全球新能源汽车市场持续高速增长，中国作为最大单一市场，占据全球55%的份额。"
        "政策驱动向市场驱动转型完成，消费者自发购买意愿显著增强。"
    )
    doc.add_paragraph(
        "核心数据：市场规模1.2万亿元 | 同比增长37.5% | 渗透率40.3% | 出口量120万辆"
    )

    doc.add_heading("二、市场规模与增长", level=1)
    doc.add_paragraph(
        "国内市场：2025年销量950万辆，预计2026年突破1200万辆，CAGR达28%。"
        "出口市场：2025年出口120万辆，同比增长50%，欧洲和东南亚为主要目的地。"
        "充电基础设施：全国充电桩保有量突破800万个，车桩比降至1.2:1。"
    )

    table = doc.add_table(rows=5, cols=4)
    table.style = 'Table Grid'
    headers = ["年份", "销量(万辆)", "同比增长", "渗透率"]
    for i, h in enumerate(headers):
        table.rows[0].cells[i].text = h
    data = [
        ["2022", "688", "93%", "26%"],
        ["2023", "780", "13%", "31%"],
        ["2024", "860", "10%", "35%"],
        ["2025", "950", "37.5%", "40%"],
    ]
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            table.rows[r + 1].cells[c].text = val

    doc.add_heading("三、竞争格局", level=1)
    doc.add_paragraph(
        "比亚迪以35%的市场份额稳居榜首，全年销量超330万辆。"
        "特斯拉中国市场份额降至7%，面临国产品牌强力竞争。"
        "新势力阵营分化加剧：理想汽车盈利突破，蔚来换电模式获认可，小鹏智驾领先。"
        "华为赋能模式崛起：问界、智界、享界系列快速放量。"
    )

    table2 = doc.add_table(rows=6, cols=4)
    table2.style = 'Table Grid'
    headers2 = ["品牌", "2025销量(万辆)", "市场份额", "同比变化"]
    for i, h in enumerate(headers2):
        table2.rows[0].cells[i].text = h
    data2 = [
        ["比亚迪", "330", "35%", "+5pp"],
        ["特斯拉中国", "66", "7%", "-3pp"],
        ["吉利新能源", "85", "9%", "+2pp"],
        ["长安新能源", "55", "6%", "+1pp"],
        ["理想汽车", "50", "5%", "+2pp"],
    ]
    for r, row_data in enumerate(data2):
        for c, val in enumerate(row_data):
            table2.rows[r + 1].cells[c].text = val

    doc.add_heading("四、技术趋势", level=1)
    doc.add_paragraph(
        "电池技术：固态电池进入量产前夜，宁德时代、比亚迪、卫蓝新能源竞相突破。"
        "能量密度突破400Wh/kg，充电速度实现10分钟充至80%。"
    )
    doc.add_paragraph(
        "智能驾驶：L3级自动驾驶法规落地，城市NOA全面开城。"
        "华为ADS 3.0、小鹏XNGP、理想AD Max三强争霸。"
    )
    doc.add_paragraph(
        "车路云一体化：5G-V2X基础设施加速部署，智慧道路覆盖超10万公里。"
    )

    doc.add_heading("五、投资建议", level=1)
    doc.add_paragraph(
        "建议关注：1) 电池产业链龙头——宁德时代、比亚迪电池、亿纬锂能；"
        "2) 智能驾驶核心标的——德赛西威、伯特利、经纬恒润；"
        "3) 充电桩运营——特锐德、星星充电。"
    )
    doc.add_paragraph(
        "风险提示：1) 补贴退坡超预期；2) 原材料价格波动；3) 地缘政治影响出口。"
    )

    out_path = Path("data/uploads/sample_nev_report.docx")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out_path))
    print("[1] Sample DOCX created: %s (%dKB)" % (out_path, out_path.stat().st_size // 1024))
    return str(out_path)


def enrich_slide_data(slide_data_list, extraction):
    """Post-process: attach tables to corresponding slides and add rich slide types."""
    tables = extraction.tables

    section_table_map = {}
    if len(tables) >= 1:
        section_table_map[1] = tables[0]
    if len(tables) >= 2:
        section_table_map[2] = tables[1]

    for idx, table_data in section_table_map.items():
        slide_idx = idx + 1
        if slide_idx < len(slide_data_list):
            slide_data_list[slide_idx]["table_data"] = table_data
            slide_data_list[slide_idx]["slide_type"] = "data"

    kpi_slide = {
        "slide_type": "findings",
        "title": "核心KPI指标",
        "content": "",
        "items": [
            "950万辆  同比+37.5%  2025年销量",
            "1.2万亿元  市场规模",
            "40.3%  渗透率突破",
            "120万辆  出口量  同比+50%",
            "800万个  充电桩保有量",
        ],
        "table_data": [],
        "extra_tables": [],
        "images": [],
        "source_text": "",
        "section_number": 0,
        "section_summary": "",
        "insight_text": "",
        "kpi_data": [],
        "comparison_data": [],
    }
    slide_data_list.insert(2, kpi_slide)

    toc_slide = {
        "slide_type": "toc",
        "title": "报告目录",
        "content": "",
        "items": [
            "行业概览与宏观环境",
            "市场规模与增长趋势",
            "竞争格局深度分析",
            "技术趋势与突破",
            "投资建议与风险提示",
        ],
        "table_data": [],
        "extra_tables": [],
        "images": [],
        "source_text": "",
        "section_number": 0,
        "section_summary": "",
        "insight_text": "",
        "kpi_data": [],
        "comparison_data": [],
    }
    slide_data_list.insert(1, toc_slide)

    return slide_data_list


def main():
    print("=" * 60)
    print("  Data-Driven PPT Generation Pipeline - Live Demo")
    print("=" * 60)

    docx_path = create_sample_docx()

    print("\n[2] Extracting data from DOCX...")
    adapter = PptInputAdapter()
    extraction = adapter.extract([docx_path])
    print("    Title: %s" % extraction.title)
    print("    Sections: %d" % len(extraction.sections))
    print("    Tables: %d" % len(extraction.tables))
    print("    Key topics: %s" % extraction.key_topics)
    for i, s in enumerate(extraction.sections):
        print("    Section %d: %s (%d chars, %d points)" % (i, s.title[:40], len(s.content), len(s.points)))

    print("\n[3] Extracting PPT requirements...")
    extractor = PptRequirementExtractor()
    requirement = extractor.extract(extraction, "请基于这份材料生成一份关于新能源汽车的PPT报告")
    print("    Topic: %s" % requirement.topic)
    print("    Focus: %s" % requirement.focus)
    print("    Page count: %d" % requirement.page_count)

    print("\n[4] Checking data gaps...")
    supplementer = PptDataSupplementer()
    gaps = supplementer.analyze_gaps(extraction, requirement)
    print("    Gaps found: %d" % len(gaps))

    print("\n[5] Building slide data...")
    builder = SlideDataBuilder()
    slide_data_list = builder.build_list(
        extraction.sections,
        add_cover=True,
        add_end=True,
        title="2026年中国新能源汽车产业深度研究报告",
    )
    print("    Initial slides: %d" % len(slide_data_list))

    print("\n[5b] Enriching slide data (tables, KPI, TOC)...")
    slide_data_list = enrich_slide_data(slide_data_list, extraction)
    print("    Enriched slides: %d" % len(slide_data_list))
    for i, sd in enumerate(slide_data_list):
        slide_type = sd.get("slide_type", "?")
        title = sd.get("title", "?")[:30]
        items_count = len(sd.get("items", []))
        table_count = len(sd.get("table_data", []))
        print("    Slide %d: [%s] %s (items=%d, tables=%d)" % (i, slide_type, title, items_count, table_count))

    print("\n[6] Generating PPTX...")
    output_dir = Path("data/reports")
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = str(output_dir / "data_driven_nev_report.pptx")

    editor = PptStructureEditor()
    result = editor.edit(slide_data_list, pptx="dummy", output_path=output_path)

    if result and result.success:
        file_size = Path(output_path).stat().st_size
        print("    SUCCESS! PPTX: %s (%dKB)" % (output_path, file_size // 1024))
    else:
        print("    FAILED: %s" % result)
        return

    print("\n[7] Persisting slide data...")
    task_id = "demo_nev_001"
    data_dir = Path("data/slide_data")
    data_dir.mkdir(parents=True, exist_ok=True)
    store = SlideDataStore(data_dir=str(data_dir), task_id=task_id)
    store.persist(task_id, slide_data_list)
    store.set_pptx_path(task_id, output_path)
    print("    Saved for task: %s" % task_id)

    print("\n[8] Creating version snapshot...")
    revisions_dir = Path("data/revisions")
    revisions_dir.mkdir(parents=True, exist_ok=True)
    version_mgr = PptVersionManager(revisions_dir=str(revisions_dir))
    version = version_mgr.create_snapshot(task_id, output_path, "L0", "initial generation")
    print("    Snapshot: version %d" % version)

    print("\n[9] Inspecting generated PPTX...")
    from pptx import Presentation
    prs = Presentation(output_path)
    print("    Slide count: %d" % len(prs.slides))
    print("    Slide size: %.1f x %.1f inches" % (prs.slide_width/914400, prs.slide_height/914400))
    for i, slide in enumerate(prs.slides):
        texts = []
        has_table = False
        has_image = False
        for sh in slide.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip().replace("\n", " | ")[:80]
                if t:
                    texts.append(t)
            if sh.has_table:
                has_table = True
                tbl = sh.table
                texts.append("[TABLE %dx%d]" % (len(tbl.rows), len(tbl.columns)))
            if hasattr(sh, 'image'):
                try:
                    _ = sh.image
                    has_image = True
                except:
                    pass
        extras = []
        if has_table:
            extras.append("TABLE")
        if has_image:
            extras.append("IMAGE")
        extra_str = " [" + ",".join(extras) + "]" if extras else ""
        summary = " | ".join(texts[:3])
        try:
            print("    Slide %d: %s%s" % (i+1, summary, extra_str))
        except UnicodeEncodeError:
            print("    Slide %d: %s%s" % (i+1, summary.encode('ascii', 'replace').decode(), extra_str))

    print("\n" + "=" * 60)
    print("  Done! PPTX saved to: %s" % output_path)
    print("=" * 60)


if __name__ == "__main__":
    main()
